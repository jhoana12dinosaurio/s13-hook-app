import os
from pptx import Presentation
from utils.pptx_utils import (
    cargar_json,
    obtener_layout_por_nombre,
    reemplazar_marcadores,
    crear_diapositiva_con_marcadores,
    verificar_contenido_diapositivas
)

# --- INICIO DE LA SOLUCIÓN PROPUESTA ---
RUTA_PROYECTO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
RUTA_PLANTILLAS = os.path.join(RUTA_PROYECTO, 'plantilla')
RUTA_JSON = os.path.join(RUTA_PROYECTO, 'json')
RUTA_PRESENTACIONES = os.path.join(RUTA_PROYECTO, 'presentaciones')
# --- FIN DE LA SOLUCIÓN PROPUESTA ---



def generar_presentacion_semana7_1_javascript():
    """
    Genera una presentación completa para la semana 7, sesión 1 del curso de 
    JavaScript Avanzado.
    Tema: Angular: concepto, instalación, creación de proyectos, directivas y Bootstrap.
    """
    print("Iniciando generación de presentación Semana 7 - Sesión 1 (JavaScript Avanzado - Angular Intro)")
    
    # Cargar la plantilla y los datos del curso
    ruta_plantilla = os.path.join(RUTA_PLANTILLAS, 'plantilla_2.pptx')
    ruta_json = os.path.join(RUTA_JSON, 'JavaScript.json')

    if not os.path.exists(ruta_plantilla):
        print(f"Error: No se encontró la plantilla en {ruta_plantilla}")
        return
    
    if not os.path.exists(ruta_json):
        print(f"Error: No se encontró el archivo JSON en {ruta_json}")
        return

    prs = Presentation(ruta_plantilla)
    datos_curso_completo = cargar_json(ruta_json) 
    if not datos_curso_completo or 'silabo' not in datos_curso_completo: 
        print("Error: No se pudieron cargar los datos del curso o estructura JSON inválida.")
        return
    # Acceder a los datos específicos del curso dentro de la estructura
    datos_curso = datos_curso_completo['silabo']
    datos_generales = datos_curso.get('datosGenerales', {})
    docente_nombre = datos_curso.get('curso', {}).get('docente', 'Docente Desconocido') # Tomando del nivel superior si es necesario
    if docente_nombre == 'Docente Desconocido': # Fallback si no está en el nivel superior
         docente_nombre = "Ms. Johan Max A. López Heredia" # Asegurando el nombre correcto del docente

    # --- Diapositiva de Portada ---
    slide_portada = crear_diapositiva_con_marcadores(prs, "Layout_Portada_UTP_1")
    # Título simplificado para la portada
    titulo_portada = "S07.s1 - Introducción a Angular y Bootstrap" 
    reemplazar_marcadores(slide_portada, {
        "TITULO_SEMANA": titulo_portada,
        "NOMBRE_SEMANA": f"{datos_curso.get('curso', 'Curso Desconocido')} - {datos_curso.get('codigo', 'Código Desconocido')}",
        "DOCENTE_NOMBRE": docente_nombre
    })
    
    # --- Momento 1: Inicio ---
    slide_inicio = crear_diapositiva_con_marcadores(prs, "Layout_Inicio")
    reemplazar_marcadores(slide_inicio, {
        "TITULO_INICIO": "Construyendo Interfaces Web Modernas con Angular",
        "ACTIVIDAD_ROMPEHIELO": (
            "Actividad inicial (5 minutos):\n\n"
            "Piensa en una aplicación web que uses mucho (Gmail, Facebook, Netflix, etc.):\n\n"
            "- ¿Notas cómo algunas partes de la página se actualizan sin recargar toda la página?\n"
            "- ¿Cómo crees que logran esa interactividad y rapidez?\n"
            "- ¿Te imaginas construir algo así solo con HTML, CSS y JavaScript básico?\n\n"
            "Hoy conoceremos Angular, un 'framework' que nos ayuda a construir este tipo de aplicaciones complejas de forma organizada."
        )
    })
    
    # --- Momento 2: Utilidad ---
    # Diapositiva de Logro
    slide_logro = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_Logro")
    # Logro basado en la Unidad 2, enfocado en S7.s1
    logro_sesion = "Al finalizar la sesión, el estudiante identifica los conceptos básicos de Angular, configura su entorno de desarrollo inicial, crea un proyecto simple y reconoce las directivas estructurales y de atributos para manipular el DOM."
    
    reemplazar_marcadores(slide_logro, {
        "ENCABEZADO_LOGRO": "Logro de Aprendizaje de la Sesión",
        "LOGRO_TEXTO": logro_sesion
    })
    
    # Diapositiva de Utilidad y Dudas
    slide_dudas = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_dudas")
    reemplazar_marcadores(slide_dudas, {
        "ENCABEZADO_DUDAS": "Dudas Frecuentes / Repaso",
        "DUDAS_TEXTO": (
            "- ¿Qué es un 'framework'? ¿Por qué necesito uno?\n"
            "- ¿Angular es lo mismo que AngularJS? (¡No!)\n"
            "- ¿Necesito saber TypeScript? (Sí, es la base de Angular).\n"
            "- ¿Es difícil instalar y empezar con Angular?\n"
            "- ¿Qué son las 'directivas'?"
        ),
        "ENCABEZADO_IMPORTANCIA_TEMA": "Importancia del Tema",
        "IMPORTANCIA_TEXTO": (
            "Aprender Angular es relevante porque:\n\n"
            "• Es uno de los frameworks frontend más populares y demandados en la industria.\n"
            "• Permite crear aplicaciones web complejas (Single Page Applications - SPAs) de forma estructurada.\n"
            "• Fomenta buenas prácticas de desarrollo (componentes, modularidad).\n"
            "• Es mantenido por Google, asegurando soporte y evolución.\n"
            "• Facilita la creación de interfaces de usuario interactivas y dinámicas."
        )
    })
    
    # Diapositiva de Saberes Previos
    slide_saberes = crear_diapositiva_con_marcadores(prs, "Layout_Saberes")
    reemplazar_marcadores(slide_saberes, {
        "TITULO_CONOCIMIENTOS": "Recordando JavaScript y la Web",
        "DINAMICA_PREVIA": (
            "Recordemos brevemente (5 minutos):\n\n"
            "1. ¿Qué lenguaje usamos para dar comportamiento e interactividad a las páginas web? (JavaScript).\n\n"
            "2. ¿Qué usamos para la estructura? (HTML). ¿Y para el estilo? (CSS).\n\n"
            "3. ¿Qué era NodeJS? (Un entorno para ejecutar JavaScript fuera del navegador, en el servidor).\n\n"
            "4. ¿Qué es TypeScript (mencionado en temario NodeJS)? (Un 'superset' de JavaScript que añade tipos estáticos).\n\n"
            "Angular usa TypeScript y nos ayuda a organizar HTML, CSS y la lógica JS/TS."
        )
    })
    
    # --- Momento 3: Transformación ---
    # Diapositiva 1: ¿Qué es Angular?
    slide_desarrollo1 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo1, {
        "SUBTEMA_TITULO": "¿Qué es Angular?",
        "CONTENIDO_TEXTO": (
            "Angular es una plataforma y framework de desarrollo frontend, basado en TypeScript, liderado por Google.\n\n"
            "• Permite construir aplicaciones web del lado del cliente (lo que ve el usuario en el navegador).\n"
            "• Enfocado en crear Single Page Applications (SPAs): Apps que cargan una sola vez y actualizan su contenido dinámicamente.\n"
            "• Basado en Componentes: La interfaz se divide en piezas reutilizables (como bloques de LEGO).\n"
            "• Utiliza TypeScript: Un lenguaje que añade tipos a JavaScript, ayudando a prevenir errores.\n\n"
            "No confundir con AngularJS (versión 1.x), que es muy diferente."
        )
    })
    
    # Diapositiva 2: Instalación y Creación de Proyecto
    slide_desarrollo2 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo2, {
        "SUBTEMA_TITULO_CODIGO": "Instalación y Creación de Proyecto (CLI)",
        # Comandos básicos, no código complejo
        "EJEMPLO_CODIGO": (
            "# Se necesita NodeJS y npm (Node Package Manager) instalados previamente.\n\n"
            "# 1. Instalar Angular CLI (Command Line Interface) globalmente:\n"
            "# (Se hace una sola vez en tu máquina)\n"
            "npm install -g @angular/cli\n\n"
            "# 2. Crear un nuevo proyecto Angular:\n"
            "# (Navega a la carpeta donde quieres crearlo)\n"
            "ng new mi-primera-app-angular\n"
            "# (Te hará preguntas sobre routing, CSS, etc. Puedes aceptar defaults).\n\n"
            "# 3. Entrar a la carpeta del proyecto:\n"
            "cd mi-primera-app-angular\n\n"
            "# 4. Ejecutar la aplicación en modo desarrollo:\n"
            "ng serve -o\n"
            "# (-o abre automáticamente el navegador en http://localhost:4200)"
        )
    })

    # Diapositiva 3: Directivas: Modificando el DOM
    slide_desarrollo3 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo3, {
        "SUBTEMA_TITULO": "Directivas: Modificando la Apariencia y Estructura",
        "CONTENIDO_TEXTO": (
            "Las directivas son instrucciones especiales en el HTML que le dicen a Angular cómo modificar el DOM (la estructura de la página).\n\n"
            "Hay dos tipos principales:\n"
            "1. Directivas de Atributos:\n"
            "   - Cambian la apariencia o comportamiento de un elemento existente.\n"
            "   - Ejemplos comunes: `[ngClass]`, `[ngStyle]` (modifican clases o estilos CSS).\n\n"
            "2. Directivas Estructurales:\n"
            "   - Cambian la estructura del DOM añadiendo o quitando elementos.\n"
            "   - Se reconocen por el asterisco `*`.\n"
            "   - Ejemplos comunes: `*ngIf`, `*ngFor`."
        )
    })

    # Diapositiva 4: Directivas Estructurales: *ngIf y *ngFor
    slide_desarrollo4 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo4, {
        "SUBTEMA_TITULO_CODIGO": "Directivas Estructurales Clave: *ngIf y *ngFor",
        # Pseudo-HTML con directivas
        "EJEMPLO_CODIGO": (
            "# *ngIf: Muestra u oculta un elemento según una condición.\n"
            "<p *ngIf=\"usuarioLogueado\">\n"
            "  Bienvenido, usuario!\n"
            "</p>\n"
            "<p *ngIf=\"!usuarioLogueado\">\n"
            "  Por favor, inicia sesión.\n"
            "</p>\n"
            "# (usuarioLogueado sería una variable booleana en el componente TS)\n\n"
            "# *ngFor: Repite un elemento por cada ítem en una lista.\n"
            "<ul>\n"
            "  <li *ngFor=\"let producto of listaProductos\">\n"
            "    {{ producto.nombre }} - S/ {{ producto.precio }}\n"
            "  </li>\n"
            "</ul>\n"
            "# (listaProductos sería un arreglo en el componente TS)"
        )
    })

    # Diapositiva 5: Integración con Bootstrap
    slide_desarrollo5 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo5, {
        "SUBTEMA_TITULO_CODIGO": "Dando Estilo Fácil: Integración con Bootstrap",
        # Comandos de instalación
        "EJEMPLO_CODIGO": (
            "# Bootstrap es un framework CSS muy popular para crear interfaces atractivas rápidamente.\n\n"
            "# Para usar Bootstrap en Angular:\n\n"
            "# 1. Instalar Bootstrap vía npm:\n"
            "npm install bootstrap\n\n"
            "# 2. Importar los estilos CSS de Bootstrap en tu archivo \n#    `angular.json` (dentro de projects > [tu-app] > architect > build > options > styles):\n"
            # "..."
            # "\"styles\": [\n"
            # "  \"node_modules/bootstrap/dist/css/bootstrap.min.css\",\n"
            # "  \"src/styles.css\"\n"
            # "],\n"
            # "..."
            
            "# 3. (Opcional) Instalar dependencias JS si necesitas componentes interactivos:\n"
            "# npm install @popperjs/core\n"
            "# Y añadir a 'scripts' en angular.json.\n\n"
            "# Una vez configurado, puedes usar las clases CSS de Bootstrap \n# en tus plantillas HTML directamente."
        )
    })
    
    # --- Momento 4: Práctica ---
    slide_practica = crear_diapositiva_con_marcadores(prs, "Layout_Practica")
    reemplazar_marcadores(slide_practica, {
        "TITULO_PRACTICA": "Práctica: Creando Nuestro Primer Proyecto Angular",
        "INSTRUCCIONES_PRACTICA": (
             "Actividad Práctica Guiada (Sencilla):\n\n"
             "Objetivo: Crear un proyecto Angular básico y ver su estructura.\n\n"
             "Pasos:\n"
             "1. Asegúrate de tener NodeJS y npm instalados.\n"
             "2. Abre una terminal o línea de comandos.\n"
             "3. Si no tienes Angular CLI, instálalo: `npm install -g @angular/cli`\n"
             "4. Navega a una carpeta donde guardarás tus proyectos.\n"
             "5. Crea un nuevo proyecto:\n"
             "   `ng new mi-primera-app` (Acepta las opciones por defecto).\n"
             "6. Espera a que se instalen las dependencias (puede tardar un poco).\n"
             "7. Entra a la carpeta del proyecto:\n"
             "   `cd mi-primera-app`\n"
             "8. Ejecuta la aplicación:\n"
             "   `ng serve -o`\n"
             "9. Observa la página que se abre en el navegador.\n"
             "10. (Opcional) Abre la carpeta del proyecto en VS Code u otro editor y explora brevemente la estructura de archivos creada (src/app, etc.).\n\n"
             "¡El objetivo principal es lograr crear y ejecutar el proyecto base!"
        ),
        "RECURSOS_NECESARIOS": (
             "• NodeJS y npm instalados.\n"
             "• Terminal o línea de comandos.\n"
             "• Navegador web.\n"
             "• (Opcional) Editor de código como VS Code.\n"
             "• Tiempo estimado: 25-30 minutos (depende de la velocidad de instalación)."
        )
    })
    
    # --- Momento 5: Cierre ---
    # Diapositiva de Cierre y Síntesis
    slide_cierre = crear_diapositiva_con_marcadores(prs, "Layout_Cierre")
    reemplazar_marcadores(slide_cierre, {
        "TITULO_CIERRE": "Resumen: Iniciando con Angular",
        "RESUMEN_IDEAS": (
             "En esta sesión hemos:\n\n"
             "• Introducido Angular como un framework frontend basado en TypeScript.\n"
             "• Visto los pasos para instalar Angular CLI y crear un nuevo proyecto.\n"
             "• Comprendido el concepto de Directivas (de Atributos y Estructurales).\n"
             "• Identificado las directivas clave `*ngIf` y `*ngFor`.\n"
             "• Conocido cómo integrar Bootstrap para mejorar el estilo visual.\n"
             "• Realizado una práctica para crear y ejecutar un proyecto Angular básico."
        ),
        "PREGUNTAS_METACOGNITIVAS": (
             "Reflexionemos:\n\n"
             "1. ¿Qué ventaja principal le ves a usar un framework como Angular en lugar de solo HTML/CSS/JS?\n"
             "2. ¿Para qué tipo de aplicaciones web crees que Angular sería más útil?\n"
             "3. ¿Qué te pareció el proceso de crear y ejecutar el proyecto inicial?"
        )
    })
    
    # Diapositiva de Preguntas
    slide_preguntas = crear_diapositiva_con_marcadores(prs, "Layout_Preguntas")
    reemplazar_marcadores(slide_preguntas, {
        "TITULO_PREGUNTAS": "¿Dudas sobre Angular, CLI, Directivas o Bootstrap?",
        "DUDAS_LISTADO": (
             "Ahora es el momento de resolver dudas sobre:\n\n"
             "• Qué es Angular y su propósito.\n"
             "• La instalación de Angular CLI (`npm install`).\n"
             "• Los comandos `ng new` y `ng serve`.\n"
             "• El concepto de directivas (`*ngIf`, `*ngFor`).\n"
             "• Cómo añadir Bootstrap al proyecto.\n"
             "• Cualquier problema encontrado durante la práctica.\n\n"
             "¡Asegurémonos de tener una base sólida!"
        )
    })
    
    # Diapositiva Final de Agradecimiento
    slide_final = crear_diapositiva_con_marcadores(prs, "1_Layout_Preguntas")
    reemplazar_marcadores(slide_final, {
        "TITULO_AGRADECIMIENTO": "¡Gracias por dar el primer paso con Angular!\n\nEn la próxima sesión, exploraremos los elementos fundamentales de Angular: Componentes, Navegación y Pipes."
    })

    # --- FIN DE LAS DIAPOSITIVAS ---

    # Verificar contenido final (opcional pero útil para depurar)
    # verificar_contenido_diapositivas(prs)

    # Crear directorio si no existe
    os.makedirs(RUTA_PRESENTACIONES, exist_ok=True)

    # Guardar la presentación generada
    nombre_archivo = 'presentacion_S07_s1_introduccion_angular_bootstrap.pptx'
    ruta_salida = os.path.join(RUTA_PRESENTACIONES, nombre_archivo)
    try:
        prs.save(ruta_salida)
        print(f"Presentación de la Semana 7 - Sesión 1 (JavaScript Avanzado) generada exitosamente en: {ruta_salida}")
    except Exception as e:
        print(f"Error al guardar la presentación: {e}")

# --- Fin de la función ---

def generar_presentacion_semana8_1_javascript_x():
    """
    Genera una presentación completa para la semana 8, sesión 1 del curso de 
    JavaScript Avanzado.
    Tema: Elementos de Angular: navegación, componentes.
    """
    print("Iniciando generación de presentación Semana 8 - Sesión 1 (JavaScript Avanzado - Elementos de Angular)")
    
    # Cargar la plantilla y los datos del curso
    ruta_plantilla = os.path.join(RUTA_PLANTILLAS, 'plantilla_2.pptx')
    ruta_json = os.path.join(RUTA_JSON, 'JavaScript.json')

    if not os.path.exists(ruta_plantilla):
        print(f"Error: No se encontró la plantilla en {ruta_plantilla}")
        return
    
    if not os.path.exists(ruta_json):
        print(f"Error: No se encontró el archivo JSON en {ruta_json}")
        return

    prs = Presentation(ruta_plantilla)
    datos_curso_completo = cargar_json(ruta_json) 
    if not datos_curso_completo or 'silabo' not in datos_curso_completo: 
        print("Error: No se pudieron cargar los datos del curso o estructura JSON inválida.")
        return
    
    datos_curso = datos_curso_completo['silabo']
    docente_nombre = "Ms. Johan Max A. López Heredia"

    # --- Diapositiva de Portada ---
    slide_portada = crear_diapositiva_con_marcadores(prs, "Layout_Portada_UTP_1")
    titulo_portada = "S08.s1 - Elementos de Angular: Navegación y Componentes"
    reemplazar_marcadores(slide_portada, {
        "TITULO_SEMANA": titulo_portada,
        "NOMBRE_SEMANA": f"{datos_curso.get('curso', 'Curso Desconocido')} - {datos_curso.get('codigo', 'Código Desconocido')}",
        "DOCENTE_NOMBRE": docente_nombre
    })
    
    # --- Momento 1: Inicio ---
    slide_inicio = crear_diapositiva_con_marcadores(prs, "Layout_Inicio")
    reemplazar_marcadores(slide_inicio, {
        "TITULO_INICIO": "Construyendo Aplicaciones Web Dinámicas con Componentes y Navegación Angular",
        "ACTIVIDAD_ROMPEHIELO": (
            "Actividad inicial (5 minutos):\n\n"
            "Pensemos en una aplicación web como Netflix o Spotify:\n\n"
            "- ¿Cómo crees que se organiza la navegación entre diferentes secciones?\n"
            "- ¿Qué elementos comunes ves en diferentes páginas (cabeceras, menús, etc.)?\n"
            "- ¿Has notado cómo al navegar por estas apps, algunas partes cambian mientras otras permanecen iguales?\n\n"
            "Hoy aprenderemos cómo Angular organiza la estructura de una aplicación a través de componentes y cómo implementa la navegación entre diferentes vistas."
        )
    })
    
    # --- Momento 2: Utilidad ---
    # Diapositiva de Logro
    slide_logro = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_Logro")
    logro_sesion = "Al finalizar la sesión, el estudiante diseña y desarrolla aplicaciones Angular con múltiples vistas, implementando un sistema de navegación eficiente y creando componentes reutilizables que mejoran la experiencia del usuario y simplifican el mantenimiento del código."
    
    reemplazar_marcadores(slide_logro, {
        "ENCABEZADO_LOGRO": "Logro de Aprendizaje de la Sesión",
        "LOGRO_TEXTO": logro_sesion
    })
    
    # Diapositiva de Utilidad y Dudas
    slide_dudas = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_dudas")
    reemplazar_marcadores(slide_dudas, {
        "ENCABEZADO_DUDAS": "Dudas Frecuentes / Repaso",
        "DUDAS_TEXTO": (
            "- ¿Qué son los componentes en Angular y por qué son importantes?\n"
            "- ¿Cómo se navega entre diferentes vistas en una SPA?\n"
            "- ¿Qué es Angular Router y cómo funciona?\n"
            "- ¿Qué son las rutas parametrizadas?\n"
            "- ¿Cómo puedo reutilizar elementos comunes en diferentes páginas?"
        ),
        "ENCABEZADO_IMPORTANCIA_TEMA": "Importancia del Tema",
        "IMPORTANCIA_TEXTO": (
            "Dominar componentes y navegación en Angular es fundamental porque:\n\n"
            "• Permite crear aplicaciones web modernas de múltiples vistas sin recargar la página.\n"
            "• La arquitectura basada en componentes promueve la reutilización de código y facilita el mantenimiento.\n"
            "• Mejora la experiencia del usuario con navegación fluida y tiempos de carga reducidos.\n"
            "• Permite desarrollar interfaces complejas dividiéndolas en partes más pequeñas y manejables.\n"
            "• Es la base para crear aplicaciones web escalables y profesionales."
        )
    })
    
    # Diapositiva de Saberes Previos
    slide_saberes = crear_diapositiva_con_marcadores(prs, "Layout_Saberes")
    reemplazar_marcadores(slide_saberes, {
        "TITULO_CONOCIMIENTOS": "Recordando Conceptos Fundamentales de Angular",
        "DINAMICA_PREVIA": (
            "Recordemos brevemente (5 minutos):\n\n"
            "1. ¿Qué es Angular y para qué tipos de aplicaciones está diseñado? (Framework para SPAs).\n\n"
            "2. ¿Qué herramienta usamos para crear un proyecto Angular? (Angular CLI).\n\n"
            "3. ¿Qué son las directivas en Angular y qué tipos existen? (Estructurales como *ngFor, *ngIf y de atributo como [ngClass]).\n\n"
            "4. ¿Cómo se estructura un proyecto Angular básico? (Carpetas src/app, componentes, módulos, etc.).\n\n"
            "Hoy profundizaremos en los componentes como bloques de construcción y en cómo implementar navegación entre diferentes vistas."
        )
    })
    
    # --- Momento 3: Transformación ---
    # Diapositiva 1: Componentes en Angular - Concepto
    slide_desarrollo1 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo1, {
        "SUBTEMA_TITULO": "Componentes: Los Bloques Fundamentales de Angular",
        "CONTENIDO_TEXTO": (
            "Un componente en Angular es una unidad modular y autónoma que encapsula:\n\n"
            "• Plantilla HTML: Define la estructura y apariencia (Vista)\n"
            "• Clase TypeScript: Contiene propiedades y métodos (Lógica)\n"
            "• Estilos CSS: Definen la apariencia específica del componente\n"
            "• Metadata: Configura el comportamiento mediante el decorador @Component\n\n"
            "Características clave:\n"
            "• Reutilizables: Se pueden usar en múltiples lugares\n"
            "• Anidables: Un componente puede contener otros componentes\n"
            "• Aislados: Cada componente tiene su propio ámbito\n"
            "• Testeables: Facilitan las pruebas unitarias\n\n"
            "Los componentes implementan el patrón de arquitectura Model-View-ViewModel (MVVM)"
        )
    })
    
    # Diapositiva 2: Creando Componentes
    slide_desarrollo2 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo2, {
        "SUBTEMA_TITULO_CODIGO": "Creando y Estructurando Componentes",
        "EJEMPLO_CODIGO": (
            "# 1. Crear un componente con Angular CLI\n"
            "ng generate component componentes/header\n"
            "ng generate component componentes/footer\n"
            "ng generate component paginas/inicio\n"
            "ng generate component paginas/productos\n\n"
            "# 2. Estructura de un componente (header.component.ts)\n"
            "import { Component } from '@angular/core';\n\n"
            "@Component({\n"
            "  selector: 'app-header',       // Etiqueta HTML para usar el componente\n"
            "  templateUrl: './header.component.html',  // Plantilla HTML\n"
            "  styleUrls: ['./header.component.css']    // Estilos CSS\n"
            "})\n"
            "export class HeaderComponent {\n"
            "  titulo = 'Mi Aplicación Angular';\n"
            "  mostrarMenu = true;\n\n"
            "  toggleMenu() {\n"
            "    this.mostrarMenu = !this.mostrarMenu;\n"
            "  }\n"
            "}"
        )
    })
    
    # Diapositiva 3: Uso de Componentes
    slide_desarrollo3 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo3, {
        "SUBTEMA_TITULO_CODIGO": "Utilizando Componentes en la Aplicación",
        "EJEMPLO_CODIGO": (
            "# 1. Plantilla HTML del componente (header.component.html)\n"
            "<header class=\"main-header\">\n"
            "  <h1>{{ titulo }}</h1>\n"
            "  <nav *ngIf=\"mostrarMenu\">\n"
            "    <ul>\n"
            "      <li><a routerLink=\"/\">Inicio</a></li>\n"
            "      <li><a routerLink=\"/productos\">Productos</a></li>\n"
            "      <li><a routerLink=\"/contacto\">Contacto</a></li>\n"
            "    </ul>\n"
            "  </nav>\n"
            "  <button (click)=\"toggleMenu()\">Mostrar/Ocultar Menú</button>\n"
            "</header>\n\n"
            "# 2. Usando el componente en app.component.html\n"
            "<div class=\"container\">\n"
            "  <app-header></app-header>\n"
            "  \n"
            "  <!-- El router-outlet es donde se mostrarán los componentes de ruta -->\n"
            "  <main>\n"
            "    <router-outlet></router-outlet>\n"
            "  </main>\n"
            "  \n"
            "  <app-footer></app-footer>\n"
            "</div>"
        )
    })
    
    # Diapositiva 4: Navegación en Angular - Conceptos
    slide_desarrollo4 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo4, {
        "SUBTEMA_TITULO": "Angular Router: Navegación en SPAs",
        "CONTENIDO_TEXTO": (
            "El Angular Router es un servicio que permite la navegación entre vistas en una SPA:\n\n"
            "• Navegación sin recarga de página: Actualiza solo la parte necesaria de la interfaz\n"
            "• Gestión del historial del navegador: Funcionalidad de adelante/atrás\n"
            "• Carga perezosa (Lazy Loading): Carga de módulos bajo demanda\n"
            "• Protección de rutas: Guards para controlar el acceso\n"
            "• Resolución de datos: Obtención de datos antes de mostrar componentes\n\n"
            "Elementos clave:\n"
            "• Routes: Definiciones de rutas que asocian paths a componentes\n"
            "• RouterModule: Módulo que proporciona el servicio de enrutamiento\n"
            "• router-outlet: Directiva que marca dónde se renderizan los componentes\n"
            "• routerLink: Directiva para enlaces de navegación\n"
            "• Router: Servicio para navegación programática"
        )
    })
    
    # Diapositiva 5: Implementando Navegación - Código
    slide_desarrollo5 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo5, {
        "SUBTEMA_TITULO_CODIGO": "Configurando Rutas en Angular",
        "EJEMPLO_CODIGO": (
            "# 1. Configurar rutas en app-routing.module.ts\n"
            "import { NgModule } from '@angular/core';\n"
            "import { Routes, RouterModule } from '@angular/router';\n"
            "import { InicioComponent } from './paginas/inicio/inicio.component';\n"
            "import { ProductosComponent } from './paginas/productos/productos.component';\n"
            "import { ContactoComponent } from './paginas/contacto/contacto.component';\n"
            "import { NotFoundComponent } from './paginas/not-found/not-found.component';\n\n"
            "const routes: Routes = [\n"
            "  { path: '', component: InicioComponent },  // Ruta raíz\n"
            "  { path: 'productos', component: ProductosComponent },\n"
            "  { path: 'contacto', component: ContactoComponent },\n"
            "  { path: '**', component: NotFoundComponent }  // Ruta comodín para 404\n"
            "];\n\n"
            "@NgModule({\n"
            "  imports: [RouterModule.forRoot(routes)],\n"
            "  exports: [RouterModule]\n"
            "})\n"
            "export class AppRoutingModule { }"
        )
    })
    
    # Diapositiva 6: Rutas Parametrizadas
    slide_desarrollo6 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo6, {
        "SUBTEMA_TITULO_CODIGO": "Rutas con Parámetros",
        "EJEMPLO_CODIGO": (
            "# 1. Configurar una ruta con parámetros\n"
            "const routes: Routes = [\n"
            "  ...\n"
            "  { path: 'producto/:id', component: DetalleProductoComponent },\n"
            "  ...\n"
            "];\n\n"
            "# 2. Enlazar a una ruta parametrizada\n"
            "<a [routerLink]=\"['/producto', producto.id]\">Ver detalles</a>\n\n"
            "# 3. Acceder a los parámetros en el componente\n"
            "import { ActivatedRoute } from '@angular/router';\n\n"
            "export class DetalleProductoComponent implements OnInit {\n"
            "  productoId: string;\n\n"
            "  constructor(private route: ActivatedRoute) { }\n\n"
            "  ngOnInit(): void {\n"
            "    // Acceder al parámetro 'id' de la URL\n"
            "    this.productoId = this.route.snapshot.paramMap.get('id');\n"
            "    \n"
            "    // O suscribirse a cambios en los parámetros\n"
            "    this.route.paramMap.subscribe(params => {\n"
            "      this.productoId = params.get('id');\n"
            "      // Cargar datos del producto según ID\n"
            "    });\n"
            "  }\n"
            "}"
        )
    })
    
    # Diapositiva 7: Navegación Programática
    slide_desarrollo7 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo7, {
        "SUBTEMA_TITULO_CODIGO": "Navegación Programática y Guards",
        "EJEMPLO_CODIGO": (
            "# 1. Navegación desde código TypeScript\n"
            "import { Router } from '@angular/router';\n\n"
            "export class LoginComponent {\n"
            "  constructor(private router: Router) { }\n\n"
            "  login() {\n"
            "    // Lógica de autenticación...\n"
            "    if (autenticacionExitosa) {\n"
            "      // Navegar a otra ruta\n"
            "      this.router.navigate(['/dashboard']);\n"
            "      \n"
            "      // Alternativa con parámetros de consulta\n"
            "      this.router.navigate(['/dashboard'], {\n"
            "        queryParams: { tab: 'recientes' }\n"
            "      });\n"
            "    }\n"
            "  }\n"
            "}\n\n"
            "# 2. Creando un Guard para proteger rutas\n"
            "ng generate guard guards/auth\n\n"
            "# El guard generado implementa interfaces como CanActivate\n"
            "# y permite controlar el acceso a rutas basado en condiciones"
        )
    })
    
    # --- Momento 4: Práctica ---
    slide_practica = crear_diapositiva_con_marcadores(prs, "Layout_Practica")
    reemplazar_marcadores(slide_practica, {
        "TITULO_PRACTICA": "Práctica: Creando una Mini-Aplicación con Componentes y Navegación",
        "INSTRUCCIONES_PRACTICA": (
             "Actividad Práctica Guiada:\n\n"
             "Objetivo: Crear una pequeña aplicación de catálogo de productos con navegación.\n\n"
             "Pasos:\n"
             "1. Crear un nuevo proyecto Angular o usar uno existente\n"
             "2. Generar los siguientes componentes:\n"
             "   - HeaderComponent: Para la barra de navegación\n"
             "   - HomeComponent: Página de inicio\n"
             "   - ProductListComponent: Lista de productos\n"
             "   - ProductDetailComponent: Detalle de un producto\n"
             "   - NotFoundComponent: Página 404\n"
             "3. Configurar las rutas en app-routing.module.ts:\n"
             "   - Ruta raíz ('/') para Home\n"
             "   - Ruta '/productos' para ProductList\n"
             "   - Ruta '/producto/:id' para ProductDetail\n"
             "   - Ruta '**' para NotFound\n"
             "4. Implementar la navegación:\n"
             "   - Enlaces en el HeaderComponent\n"
             "   - Lista de productos con enlaces al detalle\n"
             "   - Botón 'Volver' en el detalle de producto\n"
             "5. (Opcional) Añadir estilos con Bootstrap para mejorar la apariencia"
        ),
        "RECURSOS_NECESARIOS": (
             "• Proyecto Angular configurado\n"
             "• Conocimientos básicos de TypeScript\n"
             "• Datos de ejemplo para productos (pueden ser hardcodeados)\n"
             "• Tiempo estimado: 30-40 minutos\n"
             "• Opcional: Bootstrap instalado y configurado"
        )
    })
    
    # --- Momento 5: Cierre ---
    # Diapositiva de Cierre y Síntesis
    slide_cierre = crear_diapositiva_con_marcadores(prs, "Layout_Cierre")
    reemplazar_marcadores(slide_cierre, {
        "TITULO_CIERRE": "Resumen: Componentes y Navegación en Angular",
        "RESUMEN_IDEAS": (
             "En esta sesión hemos:\n\n"
             "• Comprendido el concepto de componentes como bloques básicos de construcción en Angular\n"
             "• Aprendido a crear componentes usando Angular CLI y a estructurarlos adecuadamente\n"
             "• Entendido cómo los componentes encapsulan HTML, lógica y estilos\n"
             "• Implementado un sistema de navegación usando Angular Router\n"
             "• Configurado rutas y asociado componentes a ellas\n"
             "• Trabajado con parámetros de ruta para pasar información entre componentes\n"
             "• Implementado navegación programática desde TypeScript"
        ),
        "PREGUNTAS_METACOGNITIVAS": (
             "Reflexionemos:\n\n"
             "1. ¿Cómo cambia la estructura de desarrollo web cuando usamos componentes en lugar de páginas tradicionales?\n\n"
             "2. ¿Qué ventajas ofrece el sistema de rutas de Angular frente a la navegación tradicional?\n\n"
             "3. ¿Qué desafíos podríamos enfrentar al desarrollar aplicaciones más complejas con múltiples niveles de navegación?"
        )
    })
    
    # Diapositiva de Preguntas
    slide_preguntas = crear_diapositiva_con_marcadores(prs, "Layout_Preguntas")
    reemplazar_marcadores(slide_preguntas, {
        "TITULO_PREGUNTAS": "¿Dudas sobre Componentes o Navegación Angular?",
        "DUDAS_LISTADO": (
             "Aclaremos dudas sobre:\n\n"
             "• El ciclo de vida de los componentes (ngOnInit, ngOnDestroy, etc.)\n"
             "• Comunicación entre componentes (Input, Output, Servicios)\n"
             "• Configuración avanzada de rutas (rutas anidadas, lazy loading)\n"
             "• Guards y protección de rutas\n"
             "• Resolvers para pre-cargar datos\n"
             "• Estrategias de navegación y manipulación de la URL\n\n"
             "¡Es el momento de resolver cualquier duda técnica!"
        )
    })
    
    # Diapositiva Final de Agradecimiento
    slide_final = crear_diapositiva_con_marcadores(prs, "1_Layout_Preguntas")
    reemplazar_marcadores(slide_final, {
        "TITULO_AGRADECIMIENTO": "¡Felicidades por construir tu primera aplicación con múltiples rutas!\n\nEn la próxima sesión, exploraremos cómo conectar nuestras aplicaciones Angular con servicios backend y APIs REST."
    })

    # --- FIN DE LAS DIAPOSITIVAS ---

    # Crear directorio si no existe
    os.makedirs(RUTA_PRESENTACIONES, exist_ok=True)

    # Guardar la presentación generada
    nombre_archivo = 'presentacion_S08_s1_elementos_angular.pptx'
    ruta_salida = os.path.join(RUTA_PRESENTACIONES, nombre_archivo)
    try:
        prs.save(ruta_salida)
        print(f"Presentación de la Semana 8 - Sesión 1 (JavaScript Avanzado - Elementos Angular) generada exitosamente en: {ruta_salida}")
    except Exception as e:
        print(f"Error al guardar la presentación: {e}")


def generar_presentacion_semana8_1_javascript():
    """
    Genera una presentación completa para la semana 8, sesión 1 del curso de
    JavaScript Avanzado.
    Tema: Elementos de Angular: navegación, componentes.
    """
    print("Iniciando generación de presentación Semana 8 - Sesión 1 (JavaScript Avanzado - Elementos de Angular)")
    
    # Cargar la plantilla y los datos del curso
    ruta_plantilla = os.path.join(RUTA_PLANTILLAS, 'plantilla_2.pptx')
    ruta_json = os.path.join(RUTA_JSON, 'JavaScript.json')

    if not os.path.exists(ruta_plantilla):
        print(f"Error: No se encontró la plantilla en {ruta_plantilla}")
        return
    
    if not os.path.exists(ruta_json):
        print(f"Error: No se encontró el archivo JSON en {ruta_json}")
        return

    prs = Presentation(ruta_plantilla)
    datos_curso_completo = cargar_json(ruta_json)
    if not datos_curso_completo or 'silabo' not in datos_curso_completo:
        print("Error: No se pudieron cargar los datos del curso o estructura JSON inválida.")
        return
    
    datos_curso = datos_curso_completo['silabo']
    docente_nombre = "Ms. Johan Max A. López Heredia"

    # --- Diapositiva de Portada ---
    slide_portada = crear_diapositiva_con_marcadores(prs, "Layout_Portada_UTP_1")
    titulo_portada = "S08.s1 - Elementos de Angular: Navegación y Componentes"
    reemplazar_marcadores(slide_portada, {
        "TITULO_SEMANA": titulo_portada,
        "NOMBRE_SEMANA": f"{datos_curso.get('curso', 'Curso Desconocido')} - {datos_curso.get('codigo', 'Código Desconocido')}",
        "DOCENTE_NOMBRE": docente_nombre
    })
    
    # --- Momento 1: Inicio ---
    slide_inicio = crear_diapositiva_con_marcadores(prs, "Layout_Inicio")
    reemplazar_marcadores(slide_inicio, {
        "TITULO_INICIO": "Construyendo Aplicaciones Web Dinámicas con Componentes y Navegación Angular",
        "ACTIVIDAD_ROMPEHIELO": (
            "Actividad inicial (5 minutos):\n\n"
            "Pensemos en una aplicación web como Netflix o Spotify:\n\n"
            "- ¿Cómo crees que se organiza la navegación entre diferentes secciones?\n"
            "- ¿Qué elementos comunes ves en diferentes páginas (cabeceras, menús, etc.)?\n"
            "- ¿Has notado cómo al navegar por estas apps, algunas partes cambian mientras otras permanecen iguales?\n\n"
            "Hoy aprenderemos cómo Angular organiza la estructura de una aplicación a través de componentes y cómo implementa la navegación entre diferentes vistas."
        )
    })
    
    # --- Momento 2: Utilidad ---
    # Diapositiva de Logro
    slide_logro = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_Logro")
    logro_sesion = "Al finalizar la sesión, el estudiante diseña y desarrolla aplicaciones Angular con múltiples vistas, implementando un sistema de navegación eficiente y creando componentes reutilizables que mejoran la experiencia del usuario y simplifican el mantenimiento del código."
    
    reemplazar_marcadores(slide_logro, {
        "ENCABEZADO_LOGRO": "Logro de Aprendizaje de la Sesión",
        "LOGRO_TEXTO": logro_sesion
    })
    
    # Diapositiva de Utilidad y Dudas
    slide_dudas = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_dudas")
    reemplazar_marcadores(slide_dudas, {
        "ENCABEZADO_DUDAS": "Dudas Frecuentes / Repaso",
        "DUDAS_TEXTO": (
            "- ¿Qué son los componentes en Angular y por qué son importantes?\n"
            "- ¿Cómo se navega entre diferentes vistas en una SPA?\n"
            "- ¿Qué es Angular Router y cómo funciona?\n"
            "- ¿Qué son las rutas parametrizadas?\n"
            "- ¿Cómo puedo reutilizar elementos comunes en diferentes páginas?"
        ),
        "ENCABEZADO_IMPORTANCIA_TEMA": "Importancia del Tema",
        "IMPORTANCIA_TEXTO": (
            "Dominar componentes y navegación en Angular es fundamental porque:\n\n"
            "• Permite crear aplicaciones web modernas de múltiples vistas sin recargar la página.\n"
            "• La arquitectura basada en componentes promueve la reutilización de código y facilita el mantenimiento.\n"
            "• Mejora la experiencia del usuario con navegación fluida y tiempos de carga reducidos.\n"
            "• Permite desarrollar interfaces complejas dividiéndolas en partes más pequeñas y manejables.\n"
            "• Es la base para crear aplicaciones web escalables y profesionales."
        )
    })
    
    # Diapositiva de Saberes Previos
    slide_saberes = crear_diapositiva_con_marcadores(prs, "Layout_Saberes")
    reemplazar_marcadores(slide_saberes, {
        "TITULO_CONOCIMIENTOS": "Recordando Conceptos Fundamentales de Angular",
        "DINAMICA_PREVIA": (
            "Recordemos brevemente (5 minutos):\n\n"
            "1. ¿Qué es Angular y para qué tipos de aplicaciones está diseñado? (Framework para SPAs).\n\n"
            "2. ¿Qué herramienta usamos para crear un proyecto Angular? (Angular CLI).\n\n"
            "3. ¿Qué son las directivas en Angular y qué tipos existen? (Estructurales como *ngFor, *ngIf y de atributo como [ngClass]).\n\n"
            "4. ¿Cómo se estructura un proyecto Angular básico? (Carpetas src/app, componentes, módulos, etc.).\n\n"
            "Hoy profundizaremos en los componentes como bloques de construcción y en cómo implementar navegación entre diferentes vistas."
        )
    })
    
    # --- Momento 3: Transformación ---
    # Diapositiva 1: Componentes en Angular - Concepto
    slide_desarrollo1 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo1, {
        "SUBTEMA_TITULO": "Componentes: Los Bloques Fundamentales de Angular",
        "CONTENIDO_TEXTO": (
            "Un componente en Angular es una unidad modular y autónoma que encapsula:\n\n"
            "• Plantilla HTML: Define la estructura y apariencia (Vista)\n"
            "• Clase TypeScript: Contiene propiedades y métodos (Lógica)\n"
            "• Estilos CSS: Definen la apariencia específica del componente\n"
            "• Metadata: Configura el comportamiento mediante el decorador @Component\n\n"
            "Características clave:\n"
            "• Reutilizables: Se pueden usar en múltiples lugares\n"
            "• Anidables: Un componente puede contener otros componentes\n"
            "• Aislados: Cada componente tiene su propio ámbito\n"
            "• Testeables: Facilitan las pruebas unitarias\n\n"
            "Los componentes implementan el patrón de arquitectura Model-View-ViewModel (MVVM)"
        )
    })
    
    # Diapositiva 2: Creando Componentes
    slide_desarrollo2 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo2, {
        "SUBTEMA_TITULO_CODIGO": "Creando y Estructurando Componentes",
        "EJEMPLO_CODIGO": (
            "# 1. Crear un componente con Angular CLI\n"
            "ng generate component componentes/header\n"
            "ng generate component componentes/footer\n"
            "ng generate component paginas/inicio\n"
            "ng generate component paginas/productos\n\n"
            "# 2. Estructura de un componente (header.component.ts)\n"
            "import { Component } from '@angular/core';\n\n"
            "@Component({\n"
            "  selector: 'app-header',       // Etiqueta HTML para usar el componente\n"
            "  templateUrl: './header.component.html',  // Plantilla HTML\n"
            "  styleUrls: ['./header.component.css']    // Estilos CSS\n"
            "})\n"
            "export class HeaderComponent {\n"
            "  titulo = 'Mi Aplicación Angular';\n"
            "  mostrarMenu = true;\n\n"
            "  toggleMenu() {\n"
            "    this.mostrarMenu = !this.mostrarMenu;\n"
            "  }\n"
            "}"
        )
    })
    
    # Diapositiva 3: Uso de Componentes
    slide_desarrollo3 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo3, {
        "SUBTEMA_TITULO_CODIGO": "Utilizando Componentes en la Aplicación",
        "EJEMPLO_CODIGO": (
            "# 1. Plantilla HTML del componente (header.component.html)\n"
            "<header class=\"main-header\">\n"
            "  <h1>{{ titulo }}</h1>\n"
            "  <nav *ngIf=\"mostrarMenu\">\n"
            "    <ul>\n"
            "      <li><a routerLink=\"/\">Inicio</a></li>\n"
            "      <li><a routerLink=\"/productos\">Productos</a></li>\n"
            "      <li><a routerLink=\"/contacto\">Contacto</a></li>\n"
            "    </ul>\n"
            "  </nav>\n"
            "  <button (click)=\"toggleMenu()\">Mostrar/Ocultar Menú</button>\n"
            "</header>\n\n"
            "# 2. Usando el componente en app.component.html\n"
            "<div class=\"container\">\n"
            "  <app-header></app-header>\n"
            "  \n"
            "  <!-- El router-outlet es donde se mostrarán los componentes de ruta -->\n"
            "  <main>\n"
            "    <router-outlet></router-outlet>\n"
            "  </main>\n"
            "  \n"
            "  <app-footer></app-footer>\n"
            "</div>"
        )
    })
    
    # Diapositiva 4: Navegación en Angular - Conceptos
    slide_desarrollo4 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo4, {
        "SUBTEMA_TITULO": "Angular Router: Navegación en SPAs",
        "CONTENIDO_TEXTO": (
            "El Angular Router es un servicio que permite la navegación entre vistas en una SPA:\n\n"
            "• Navegación sin recarga de página: Actualiza solo la parte necesaria de la interfaz\n"
            "• Gestión del historial del navegador: Funcionalidad de adelante/atrás\n"
            "• Carga perezosa (Lazy Loading): Carga de módulos bajo demanda\n"
            "• Protección de rutas: Guards para controlar el acceso\n"
            "• Resolución de datos: Obtención de datos antes de mostrar componentes\n\n"
            "Elementos clave:\n"
            "• Routes: Definiciones de rutas que asocian paths a componentes\n"
            "• RouterModule: Módulo que proporciona el servicio de enrutamiento\n"
            "• router-outlet: Directiva que marca dónde se renderizan los componentes\n"
            "• routerLink: Directiva para enlaces de navegación\n"
            "• Router: Servicio para navegación programática"
        )
    })
    
    # Diapositiva 5: Implementando Navegación - Código
    slide_desarrollo5 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo5, {
        "SUBTEMA_TITULO_CODIGO": "Configurando Rutas en Angular",
        "EJEMPLO_CODIGO": (
            "# 1. Configurar rutas en app-routing.module.ts\n"
            "import { NgModule } from '@angular/core';\n"
            "import { Routes, RouterModule } from '@angular/router';\n"
            "import { InicioComponent } from './paginas/inicio/inicio.component';\n"
            "import { ProductosComponent } from './paginas/productos/productos.component';\n"
            "import { ContactoComponent } from './paginas/contacto/contacto.component';\n"
            "import { NotFoundComponent } from './paginas/not-found/not-found.component';\n\n"
            "const routes: Routes = [\n"
            "  { path: '', component: InicioComponent },  // Ruta raíz\n"
            "  { path: 'productos', component: ProductosComponent },\n"
            "  { path: 'contacto', component: ContactoComponent },\n"
            "  { path: '**', component: NotFoundComponent }  // Ruta comodín para 404\n"
            "];\n\n"
            "@NgModule({\n"
            "  imports: [RouterModule.forRoot(routes)],\n"
            "  exports: [RouterModule]\n"
            "})\n"
            "export class AppRoutingModule { }"
        )
    })
    
    # Diapositiva 6: Rutas Parametrizadas
    slide_desarrollo6 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo6, {
        "SUBTEMA_TITULO_CODIGO": "Rutas con Parámetros",
        "EJEMPLO_CODIGO": (
            "# 1. Configurar una ruta con parámetros\n"
            "const routes: Routes = [\n"
            "  ...\n"
            "  { path: 'producto/:id', component: DetalleProductoComponent },\n"
            "  ...\n"
            "];\n\n"
            "# 2. Enlazar a una ruta parametrizada\n"
            "<a [routerLink]=\"['/producto', producto.id]\">Ver detalles</a>\n\n"
            "# 3. Acceder a los parámetros en el componente\n"
            "import { ActivatedRoute } from '@angular/router';\n\n"
            "export class DetalleProductoComponent implements OnInit {\n"
            "  productoId: string;\n\n"
            "  constructor(private route: ActivatedRoute) { }\n\n"
            "  ngOnInit(): void {\n"
            "    // Acceder al parámetro 'id' de la URL\n"
            "    this.productoId = this.route.snapshot.paramMap.get('id');\n"
            "    \n"
            "    // O suscribirse a cambios en los parámetros\n"
            "    this.route.paramMap.subscribe(params => {\n"
            "      this.productoId = params.get('id');\n"
            "      // Cargar datos del producto según ID\n"
            "    });\n"
            "  }\n"
            "}"
        )
    })
    
    # Diapositiva 7: Navegación Programática
    slide_desarrollo7 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo7, {
        "SUBTEMA_TITULO_CODIGO": "Navegación Programática y Guards",
        "EJEMPLO_CODIGO": (
            "# 1. Navegación desde código TypeScript\n"
            "import { Router } from '@angular/router';\n\n"
            "export class LoginComponent {\n"
            "  constructor(private router: Router) { }\n\n"
            "  login() {\n"
            "    // Lógica de autenticación...\n"
            "    if (autenticacionExitosa) {\n"
            "      // Navegar a otra ruta\n"
            "      this.router.navigate(['/dashboard']);\n"
            "      \n"
            "      // Alternativa con parámetros de consulta\n"
            "      this.router.navigate(['/dashboard'], {\n"
            "        queryParams: { tab: 'recientes' }\n"
            "      });\n"
            "    }\n"
            "  }\n"
            "}\n\n"
            "# 2. Creando un Guard para proteger rutas\n"
            "ng generate guard guards/auth\n\n"
            "# El guard generado implementa interfaces como CanActivate\n"
            "# y permite controlar el acceso a rutas basado en condiciones"
        )
    })
    
    # --- Momento 4: Práctica ---
    slide_practica = crear_diapositiva_con_marcadores(prs, "Layout_Practica")
    reemplazar_marcadores(slide_practica, {
        "TITULO_PRACTICA": "Práctica: Creando una Mini-Aplicación con Componentes y Navegación",
        "INSTRUCCIONES_PRACTICA": (
             "Actividad Práctica Guiada:\n\n"
             "Objetivo: Crear una pequeña aplicación de catálogo de productos con navegación.\n\n"
             "Pasos:\n"
             "1. Crear un nuevo proyecto Angular o usar uno existente\n"
             "2. Generar los siguientes componentes:\n"
             "   - HeaderComponent: Para la barra de navegación\n"
             "   - HomeComponent: Página de inicio\n"
             "   - ProductListComponent: Lista de productos\n"
             "   - ProductDetailComponent: Detalle de un producto\n"
             "   - NotFoundComponent: Página 404\n"
             "3. Configurar las rutas en app-routing.module.ts:\n"
             "   - Ruta raíz ('/') para Home\n"
             "   - Ruta '/productos' para ProductList\n"
             "   - Ruta '/producto/:id' para ProductDetail\n"
             "   - Ruta '**' para NotFound\n"
             "4. Implementar la navegación:\n"
             "   - Enlaces en el HeaderComponent\n"
             "   - Lista de productos con enlaces al detalle\n"
             "   - Botón 'Volver' en el detalle de producto\n"
             "5. (Opcional) Añadir estilos con Bootstrap para mejorar la apariencia"
        ),
        "RECURSOS_NECESARIOS": (
             "• Proyecto Angular configurado\n"
             "• Conocimientos básicos de TypeScript\n"
             "• Datos de ejemplo para productos (pueden ser hardcodeados)\n"
             "• Tiempo estimado: 30-40 minutos\n"
             "• Opcional: Bootstrap instalado y configurado"
        )
    })
    
    # --- Momento 5: Cierre ---
    # Diapositiva de Cierre y Síntesis
    slide_cierre = crear_diapositiva_con_marcadores(prs, "Layout_Cierre")
    reemplazar_marcadores(slide_cierre, {
        "TITULO_CIERRE": "Resumen: Componentes y Navegación en Angular",
        "RESUMEN_IDEAS": (
             "En esta sesión hemos:\n\n"
             "• Comprendido el concepto de componentes como bloques básicos de construcción en Angular\n"
             "• Aprendido a crear componentes usando Angular CLI y a estructurarlos adecuadamente\n"
             "• Entendido cómo los componentes encapsulan HTML, lógica y estilos\n"
             "• Implementado un sistema de navegación usando Angular Router\n"
             "• Configurado rutas y asociado componentes a ellas\n"
             "• Trabajado con parámetros de ruta para pasar información entre componentes\n"
             "• Implementado navegación programática desde TypeScript"
        ),
        "PREGUNTAS_METACOGNITIVAS": (
             "Reflexionemos:\n\n"
             "1. ¿Cómo cambia la estructura de desarrollo web cuando usamos componentes en lugar de páginas tradicionales?\n\n"
             "2. ¿Qué ventajas ofrece el sistema de rutas de Angular frente a la navegación tradicional?\n\n"
             "3. ¿Qué desafíos podríamos enfrentar al desarrollar aplicaciones más complejas con múltiples niveles de navegación?"
        )
    })
    
    # Diapositiva de Preguntas
    slide_preguntas = crear_diapositiva_con_marcadores(prs, "Layout_Preguntas")
    reemplazar_marcadores(slide_preguntas, {
        "TITULO_PREGUNTAS": "¿Dudas sobre Componentes o Navegación Angular?",
        "DUDAS_LISTADO": (
             "Aclaremos dudas sobre:\n\n"
             "• El ciclo de vida de los componentes (ngOnInit, ngOnDestroy, etc.)\n"
             "• Comunicación entre componentes (Input, Output, Servicios)\n"
             "• Configuración avanzada de rutas (rutas anidadas, lazy loading)\n"
             "• Guards y protección de rutas\n"
             "• Resolvers para pre-cargar datos\n"
             "• Estrategias de navegación y manipulación de la URL\n\n"
             "¡Es el momento de resolver cualquier duda técnica!"
        )
    })
    
    # Diapositiva Final de Agradecimiento
    slide_final = crear_diapositiva_con_marcadores(prs, "1_Layout_Preguntas")
    reemplazar_marcadores(slide_final, {
        "TITULO_AGRADECIMIENTO": "¡Felicidades por construir tu primera aplicación con múltiples rutas!\n\nEn la próxima sesión, exploraremos cómo conectar nuestras aplicaciones Angular con servicios backend y APIs REST."
    })

    # --- FIN DE LAS DIAPOSITIVAS ---

    # Crear directorio si no existe
    os.makedirs(RUTA_PRESENTACIONES, exist_ok=True)

    # Guardar la presentación generada
    nombre_archivo = 'presentacion_S08_s1_elementos_angular.pptx'
    ruta_salida = os.path.join(RUTA_PRESENTACIONES, nombre_archivo)
    try:
        prs.save(ruta_salida)
        print(f"Presentación de la Semana 8 - Sesión 1 (JavaScript Avanzado - Elementos Angular) generada exitosamente en: {ruta_salida}")
    except Exception as e:
        print(f"Error al guardar la presentación: {e}")

def generar_presentacion_semana9_1_javascript():
    """
    Genera una presentación completa para la semana 9, sesión 1 del curso de 
    JavaScript Avanzado.
    Tema: Servicios en Angular: creación de un Back-end en JSON, métodos REST, 
    herramientas de prueba, creación de servicios e implementación de una aplicación web.
    """
    print("Iniciando generación de presentación Semana 9 - Sesión 1 (JavaScript Avanzado - Servicios en Angular y REST)")
    
    # Cargar la plantilla y los datos del curso
    ruta_plantilla = os.path.join(RUTA_PLANTILLAS, 'plantilla_2.pptx')
    ruta_json = os.path.join(RUTA_JSON, 'JavaScript.json')

    if not os.path.exists(ruta_plantilla):
        print(f"Error: No se encontró la plantilla en {ruta_plantilla}")
        return
    
    if not os.path.exists(ruta_json):
        print(f"Error: No se encontró el archivo JSON en {ruta_json}")
        return

    prs = Presentation(ruta_plantilla)
    datos_curso_completo = cargar_json(ruta_json) 
    if not datos_curso_completo or 'silabo' not in datos_curso_completo: 
        print("Error: No se pudieron cargar los datos del curso o estructura JSON inválida.")
        return
    
    datos_curso = datos_curso_completo['silabo']
    docente_nombre = "Ms. Johan Max A. López Heredia"

    # --- Diapositiva de Portada ---
    slide_portada = crear_diapositiva_con_marcadores(prs, "Layout_Portada_UTP_1")
    titulo_portada = "S09.s1 - Servicios en Angular y APIs REST"
    reemplazar_marcadores(slide_portada, {
        "TITULO_SEMANA": titulo_portada,
        "NOMBRE_SEMANA": f"{datos_curso.get('curso', 'Curso Desconocido')} - {datos_curso.get('codigo', 'Código Desconocido')}",
        "DOCENTE_NOMBRE": docente_nombre
    })
    
    # --- Momento 1: Inicio ---
    slide_inicio = crear_diapositiva_con_marcadores(prs, "Layout_Inicio")
    reemplazar_marcadores(slide_inicio, {
        "TITULO_INICIO": "Conectando Angular con el Mundo: Servicios y APIs",
        "ACTIVIDAD_ROMPEHIELO": (
            "Actividad inicial (5 minutos):\n\n"
            "Piensa en aplicaciones que usas a diario como Instagram, YouTube o Spotify:\n\n"
            "- ¿De dónde crees que obtienen los datos que muestran (fotos, videos, canciones)?\n"
            "- ¿Cómo crees que se actualizan cuando hay nueva información?\n"
            "- Si publicas algo, ¿cómo crees que se guarda para que otros lo vean después?\n\n"
            "Hoy aprenderemos cómo las aplicaciones web modernas se comunican con los servidores para obtener y guardar información."
        )
    })
    
    # --- Momento 2: Utilidad ---
    slide_logro = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_Logro")
    logro_sesion = "Al finalizar la sesión, el estudiante implementa servicios en Angular para consumir APIs REST, utilizando un back-end basado en JSON y aplicando herramientas de prueba para validar el funcionamiento de la comunicación cliente-servidor."
    
    reemplazar_marcadores(slide_logro, {
        "ENCABEZADO_LOGRO": "Logro de Aprendizaje de la Sesión",
        "LOGRO_TEXTO": logro_sesion
    })
    
    slide_dudas = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_dudas")
    reemplazar_marcadores(slide_dudas, {
        "ENCABEZADO_DUDAS": "Dudas Frecuentes / Repaso",
        "DUDAS_TEXTO": (
            "- ¿Qué es una API REST y por qué es tan popular?\n"
            "- ¿Cómo se conecta Angular con servicios externos?\n"
            "- ¿Qué significa JSON Server y para qué sirve?\n"
            "- ¿Qué son los métodos HTTP (GET, POST, etc.)?\n"
            "- ¿Cómo manejo las respuestas asíncronas en Angular?"
        ),
        "ENCABEZADO_IMPORTANCIA_TEMA": "Importancia del Tema",
        "IMPORTANCIA_TEXTO": (
            "Dominar los servicios y APIs REST es fundamental porque:\n\n"
            "• Casi todas las aplicaciones web modernas requieren comunicación con servidores.\n"
            "• Es el estándar de la industria para comunicación cliente-servidor.\n"
            "• Permite separar claramente el frontend del backend (desacoplamiento).\n"
            "• Facilita el consumo de servicios externos (clima, pagos, mapas, etc.).\n"
            "• Es una habilidad muy demandada en el mercado laboral."
        )
    })
    
    # Diapositiva de Saberes Previos
    slide_saberes = crear_diapositiva_con_marcadores(prs, "Layout_Saberes")
    reemplazar_marcadores(slide_saberes, {
        "TITULO_CONOCIMIENTOS": "¿Qué Necesitamos Recordar?",
        "DINAMICA_PREVIA": (
            "Recordemos brevemente (5 minutos):\n\n"
            "1. ¿Qué es un componente en Angular? (Un bloque de construcción que combina HTML, CSS y TypeScript).\n\n"
            "2. ¿Cómo se estructura una aplicación Angular? (Mediante módulos, componentes, servicios, etc.).\n\n"
            "3. ¿Cómo podríamos compartir datos entre componentes que no tienen relación directa? (Mediante servicios).\n\n"
            "4. ¿Qué significa el término 'asíncrono' en desarrollo web? (Operaciones que ocurren fuera del flujo principal de ejecución).\n\n"
            "Hoy aprenderemos a usar servicios para la comunicación con APIs."
        )
    })
    
    # --- Momento 3: Transformación ---
    # Qué son las APIs REST
    slide_desarrollo1 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo1, {
        "SUBTEMA_TITULO": "APIs REST: La Columna Vertebral de la Web Moderna",
        "CONTENIDO_TEXTO": (
            "API REST (Representational State Transfer) es un estilo de arquitectura para sistemas web:\n\n"
            "• Basada en recursos identificados por URLs (Ej: /productos, /usuarios/1)\n"
            "• Utiliza métodos HTTP estándar:\n"
            "  - GET: Obtener datos (Consultar)\n"
            "  - POST: Crear nuevos datos\n"
            "  - PUT/PATCH: Actualizar datos existentes\n"
            "  - DELETE: Eliminar datos\n"
            "• Intercambia datos principalmente en formato JSON\n"
            "• Es stateless (sin estado): cada petición contiene toda la información necesaria\n"
            "• Soporta diferentes formatos de respuesta (JSON, XML, etc.)"
        )
    })
    
    # JSON Server
    slide_desarrollo2 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo2, {
        "SUBTEMA_TITULO_CODIGO": "JSON Server: Un Backend en Segundos",
        "EJEMPLO_CODIGO": (
            "# JSON Server es una herramienta que crea una API REST completa a partir de un archivo JSON\n\n"
            "# 1. Instalar JSON Server globalmente\n"
            "npm install -g json-server\n\n"
            "# 2. Crear un archivo db.json con datos de ejemplo\n"
            "{\n"
            "  \"productos\": [\n"
            "    { \"id\": 1, \"nombre\": \"Laptop\", \"precio\": 1200 },\n"
            "    { \"id\": 2, \"nombre\": \"Smartphone\", \"precio\": 800 }\n"
            "  ],\n"
            "  \"categorias\": [\n"
            "    { \"id\": 1, \"nombre\": \"Electrónicos\" }\n"
            "  ]\n"
            "}\n\n"
            "# 3. Iniciar el servidor\n"
            "json-server --watch db.json\n\n"
            "# ¡Listo! Ahora tienes endpoints como:\n"
            "# GET http://localhost:3000/productos\n"
            "# GET http://localhost:3000/productos/1\n"
        )
    })

    # Servicios en Angular
    slide_desarrollo3 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo3, {
        "SUBTEMA_TITULO": "Servicios en Angular: Conectando con APIs",
        "CONTENIDO_TEXTO": (
            "Los servicios en Angular son clases con un propósito específico:\n\n"
            "• Inyectables en componentes mediante Dependency Injection\n"
            "• Ideales para:\n"
            "  - Compartir datos entre componentes\n"
            "  - Encapsular lógica de negocio\n"
            "  - Comunicarse con APIs y servicios externos\n"
            "  - Centralizar la lógica de acceso a datos\n\n"
            "Para consumir APIs REST, Angular utiliza el módulo HttpClient para realizar peticiones HTTP asíncronas que devuelven Observables (parte de RxJS)."
        )
    })

    # Creando un Servicio
    slide_desarrollo4 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo4, {
        "SUBTEMA_TITULO_CODIGO": "Creando un Servicio para Consumir APIs",
        "EJEMPLO_CODIGO": (
            "# 1. Generar un servicio con Angular CLI\n"
            "ng generate service services/producto\n\n"
            "# 2. Implementar el servicio (producto.service.ts)\n"
            "import { Injectable } from '@angular/core';\n"
            "import { HttpClient } from '@angular/common/http';\n"
            "import { Observable } from 'rxjs';\n\n"
            "@Injectable({ providedIn: 'root' })\n"
            "export class ProductoService {\n"
            "  private apiUrl = 'http://localhost:3000/productos';\n\n"
            "  constructor(private http: HttpClient) { }\n\n"
            "  // Obtener todos los productos\n"
            "  getProductos(): Observable<any[]> {\n"
            "    return this.http.get<any[]>(this.apiUrl);\n"
            "  }\n\n"
            "  // Obtener un producto por ID\n"
            "  getProducto(id: number): Observable<any> {\n"
            "    return this.http.get<any>(`${this.apiUrl}/${id}`);\n"
            "  }\n"
            "}"
        )
    })

    # Métodos CRUD
    slide_desarrollo5 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo5, {
        "SUBTEMA_TITULO_CODIGO": "Implementando CRUD Completo",
        "EJEMPLO_CODIGO": (
            "# Continuando con producto.service.ts (métodos adicionales)\n\n"
            "// Crear un nuevo producto\n"
            "crearProducto(producto: any): Observable<any> {\n"
            "  return this.http.post<any>(this.apiUrl, producto);\n"
            "}\n\n"
            "// Actualizar un producto existente\n"
            "actualizarProducto(id: number, producto: any): Observable<any> {\n"
            "  return this.http.put<any>(`${this.apiUrl}/${id}`, producto);\n"
            "}\n\n"
            "// Eliminar un producto\n"
            "eliminarProducto(id: number): Observable<any> {\n"
            "  return this.http.delete<any>(`${this.apiUrl}/${id}`);\n"
            "}"
        )
    })
    
    # Uso en componentes
    slide_desarrollo6 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo6, {
        "SUBTEMA_TITULO_CODIGO": "Usando el Servicio en un Componente",
        "EJEMPLO_CODIGO": (
            "# 1. Configurar HttpClientModule en app.module.ts\n"
            "import { HttpClientModule } from '@angular/common/http';\n\n"
            "@NgModule({\n"
            "  imports: [\n"
            "    BrowserModule,\n"
            "    HttpClientModule,  // Añadir esta línea\n"
            "    // otros imports...\n"
            "  ],\n"
            "})\n\n"
            "# 2. Usar el servicio en un componente\n"
            "import { Component, OnInit } from '@angular/core';\n"
            "import { ProductoService } from '../services/producto.service';\n\n"
            "@Component({\n"
            "  selector: 'app-lista-productos',\n"
            "  templateUrl: './lista-productos.component.html'\n"
            "})\n"
            "export class ListaProductosComponent implements OnInit {\n"
            "  productos: any[] = [];\n\n"
            "  constructor(private productoService: ProductoService) { }\n\n"
            "  ngOnInit(): void {\n"
            "    this.cargarProductos();\n"
            "  }\n\n"
            "  cargarProductos(): void {\n"
            "    this.productoService.getProductos().subscribe(\n"
            "      data => this.productos = data,\n"
            "      error => console.error('Error al cargar productos:', error)\n"
            "    );\n"
            "  }\n"
            "}"
        )
    })
    
    # Herramientas de Prueba
    slide_desarrollo7 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo7, {
        "SUBTEMA_TITULO": "Herramientas para Probar APIs",
        "CONTENIDO_TEXTO": (
            "Para probar APIs REST de forma efectiva, puedes usar herramientas especializadas:\n\n"
            "• Postman: Interfaz gráfica completa para probar APIs\n"
            "  - Organiza colecciones de peticiones\n"
            "  - Guarda ejemplos de respuestas\n"
            "  - Automatiza pruebas\n\n"
            "• Insomnia: Alternativa moderna a Postman\n\n"
            "• Thunder Client: Extensión para VS Code\n\n"
            "• Browser DevTools: La pestaña Network en Chrome/Firefox\n\n"
            "• curl: Herramienta de línea de comandos para peticiones HTTP"
        )
    })
    
    # --- Momento 4: Práctica ---
    slide_practica = crear_diapositiva_con_marcadores(prs, "Layout_Practica")
    reemplazar_marcadores(slide_practica, {
        "TITULO_PRACTICA": "Práctica: App de Contactos con API REST",
        "INSTRUCCIONES_PRACTICA": (
             "Actividad Práctica Guiada (Sencilla):\n\n"
             "Objetivo: Crear una aplicación de contactos que use JSON Server.\n\n"
             "Pasos:\n"
             "1. Configurar JSON Server:\n"
             "   - Crear archivo db.json con una lista de contactos\n"
             "   - Iniciar el servidor con json-server --watch db.json\n\n"
             "2. En el proyecto Angular:\n"
             "   - Generar un servicio para contactos\n"
             "   - Implementar métodos GET, POST, PUT y DELETE\n"
             "   - Crear componente para listar contactos\n"
             "   - Crear formulario para añadir/editar contactos\n\n"
             "3. Implementar operaciones CRUD básicas:\n"
             "   - Mostrar lista de contactos\n"
             "   - Crear un nuevo contacto\n"
             "   - Actualizar un contacto existente\n"
             "   - Eliminar un contacto"
        ),
        "RECURSOS_NECESARIOS": (
             "• Proyecto Angular existente\n"
             "• JSON Server instalado (npm install -g json-server)\n"
             "• Archivo db.json creado con datos iniciales\n"
             "• HttpClientModule configurado en app.module.ts\n"
             "• Tiempo estimado: 30-40 minutos"
        )
    })
    
    # --- Momento 5: Cierre ---
    slide_cierre = crear_diapositiva_con_marcadores(prs, "Layout_Cierre")
    reemplazar_marcadores(slide_cierre, {
        "TITULO_CIERRE": "Resumen: Conectando Angular con APIs REST",
        "RESUMEN_IDEAS": (
             "En esta sesión hemos:\n\n"
             "• Entendido los conceptos fundamentales de APIs REST\n"
             "• Configurado JSON Server como un backend rápido para desarrollo\n"
             "• Creado servicios en Angular para consumir APIs\n"
             "• Implementado operaciones CRUD completas (Create, Read, Update, Delete)\n"
             "• Usado HttpClient y Observables para manejar respuestas asíncronas\n"
             "• Creado una aplicación simple pero funcional con comunicación cliente-servidor"
        ),
        "PREGUNTAS_METACOGNITIVAS": (
             "Reflexionemos:\n\n"
             "1. ¿Qué ventajas ofrece separar la lógica de acceso a datos en servicios?\n"
             "2. ¿Cómo facilita JSON Server el desarrollo de aplicaciones que necesitan un backend?\n"
             "3. ¿Por qué es importante entender el flujo asíncrono al trabajar con APIs?\n"
             "4. ¿Qué desafíos podrías enfrentar al conectar con APIs reales en producción?"
        )
    })
    
    # Diapositiva de Preguntas
    slide_preguntas = crear_diapositiva_con_marcadores(prs, "Layout_Preguntas")
    reemplazar_marcadores(slide_preguntas, {
        "TITULO_PREGUNTAS": "¿Dudas sobre Servicios o APIs REST?",
        "DUDAS_LISTADO": (
             "Resolvamos dudas sobre:\n\n"
             "• Conceptos de REST y métodos HTTP\n"
             "• Configuración y uso de JSON Server\n"
             "• Creación e inyección de servicios en Angular\n"
             "• Manejo de observables y suscripciones\n"
             "• Tratamiento de errores en peticiones HTTP\n"
             "• Buenas prácticas para estructurar servicios\n\n"
             "¡Aclaremos estos conceptos para poder avanzar!"
        )
    })
    
    # Diapositiva Final
    slide_final = crear_diapositiva_con_marcadores(prs, "1_Layout_Preguntas")
    reemplazar_marcadores(slide_final, {
        "TITULO_AGRADECIMIENTO": "¡Felicidades por dar el paso hacia aplicaciones web conectadas!\n\nEn la próxima sesión, integraremos todo lo aprendido para crear una aplicación web completa con múltiples funcionalidades."
    })

    # Guardar presentación
    os.makedirs(RUTA_PRESENTACIONES, exist_ok=True)

    nombre_archivo = 'presentacion_S09_s1_servicios_angular_rest.pptx'
    ruta_salida = os.path.join(RUTA_PRESENTACIONES, nombre_archivo)
    try:
        prs.save(ruta_salida)
        print(f"Presentación de la Semana 9 - Sesión 1 generada exitosamente en: {ruta_salida}")
    except Exception as e:
        print(f"Error al guardar la presentación: {e}")
def generar_presentacion_semana09_1_javascript_v2():
    """
    Genera una presentación completa para la semana 9, sesión 1 del curso de 
    JavaScript Avanzado, con énfasis en teoría (70/30).
    Tema: Servicios en Angular: creación de un Back-end en JSON, métodos REST, 
    herramientas de prueba, creación de servicios e implementación de una aplicación web.
    """
    print("Iniciando generación de presentación Semana 9 - Sesión 1 (v2 - Teórico) (JavaScript Avanzado - Servicios en Angular y REST)")
    
    # Cargar la plantilla y los datos del curso
    ruta_plantilla = os.path.join(RUTA_PLANTILLAS, 'plantilla_2.pptx')
    ruta_json = os.path.join(RUTA_JSON, 'JavaScript.json')

    if not os.path.exists(ruta_plantilla):
        print(f"Error: No se encontró la plantilla en {ruta_plantilla}")
        return
    
    if not os.path.exists(ruta_json):
        print(f"Error: No se encontró el archivo JSON en {ruta_json}")
        return

    prs = Presentation(ruta_plantilla)
    datos_curso_completo = cargar_json(ruta_json) 
    if not datos_curso_completo or 'silabo' not in datos_curso_completo: 
        print("Error: No se pudieron cargar los datos del curso o estructura JSON inválida.")
        return
    
    datos_curso = datos_curso_completo['silabo']
    docente_nombre = "Ms. Johan Max A. López Heredia"

    # --- Diapositiva de Portada ---
    slide_portada = crear_diapositiva_con_marcadores(prs, "Layout_Portada_UTP_1")
    titulo_portada = "S09.s1 - Servicios en Angular y APIs REST"
    reemplazar_marcadores(slide_portada, {
        "TITULO_SEMANA": titulo_portada,
        "NOMBRE_SEMANA": f"{datos_curso.get('curso', 'Curso Desconocido')} - {datos_curso.get('codigo', 'Código Desconocido')}",
        "DOCENTE_NOMBRE": docente_nombre
    })
    
    # --- Momento 1: Inicio ---
    slide_inicio = crear_diapositiva_con_marcadores(prs, "Layout_Inicio")
    reemplazar_marcadores(slide_inicio, {
        "TITULO_INICIO": "Conectando Angular con el Mundo: Servicios y APIs",
        "ACTIVIDAD_ROMPEHIELO": (
            "Actividad inicial (5 minutos):\n\n"
            "Piensa en aplicaciones que usas a diario como Instagram, YouTube o Spotify:\n\n"
            "- ¿De dónde crees que obtienen los datos que muestran (fotos, videos, canciones)?\n"
            "- ¿Cómo crees que se actualizan cuando hay nueva información?\n"
            "- Si publicas algo, ¿cómo crees que se guarda para que otros lo vean después?\n\n"
            "Hoy aprenderemos cómo las aplicaciones web modernas se comunican con los servidores para obtener y guardar información."
        )
    })
    
    # --- Momento 2: Utilidad ---
    slide_logro = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_Logro")
    logro_sesion = "Al finalizar la sesión, el estudiante implementa servicios en Angular para consumir APIs REST, utilizando un back-end basado en JSON y aplicando herramientas de prueba para validar el funcionamiento de la comunicación cliente-servidor."
    
    reemplazar_marcadores(slide_logro, {
        "ENCABEZADO_LOGRO": "Logro de Aprendizaje de la Sesión",
        "LOGRO_TEXTO": logro_sesion
    })
    
    slide_dudas = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_dudas")
    reemplazar_marcadores(slide_dudas, {
        "ENCABEZADO_DUDAS": "Dudas Frecuentes / Repaso",
        "DUDAS_TEXTO": (
            "- ¿Qué es una API REST y por qué es tan popular?\n"
            "- ¿Cómo se conecta Angular con servicios externos?\n"
            "- ¿Qué significa JSON Server y para qué sirve?\n"
            "- ¿Qué son los métodos HTTP (GET, POST, etc.)?\n"
            "- ¿Cómo manejo las respuestas asíncronas en Angular?"
        ),
        "ENCABEZADO_IMPORTANCIA_TEMA": "Importancia del Tema",
        "IMPORTANCIA_TEXTO": (
            "Dominar los servicios y APIs REST es fundamental porque:\n\n"
            "• Casi todas las aplicaciones web modernas requieren comunicación con servidores.\n"
            "• Es el estándar de la industria para comunicación cliente-servidor.\n"
            "• Permite separar claramente el frontend del backend (desacoplamiento).\n"
            "• Facilita el consumo de servicios externos (clima, pagos, mapas, etc.).\n"
            "• Es una habilidad muy demandada en el mercado laboral."
        )
    })
    
    # Diapositiva de Saberes Previos
    slide_saberes = crear_diapositiva_con_marcadores(prs, "Layout_Saberes")
    reemplazar_marcadores(slide_saberes, {
        "TITULO_CONOCIMIENTOS": "¿Qué Necesitamos Recordar?",
        "DINAMICA_PREVIA": (
            "Recordemos brevemente (5 minutos):\n\n"
            "1. ¿Qué es un componente en Angular? (Un bloque de construcción que combina HTML, CSS y TypeScript).\n\n"
            "2. ¿Cómo se estructura una aplicación Angular? (Mediante módulos, componentes, servicios, etc.).\n\n"
            "3. ¿Cómo podríamos compartir datos entre componentes que no tienen relación directa? (Mediante servicios).\n\n"
            "4. ¿Qué significa el término 'asíncrono' en desarrollo web? (Operaciones que ocurren fuera del flujo principal de ejecución).\n\n"
            "Hoy aprenderemos a usar servicios para la comunicación con APIs."
        )
    })
    
    # --- Momento 3: Transformación (Énfasis Teórico) ---
    # Slide 1: ¿Qué son las APIs REST? (Teoría)
    slide_desarrollo1 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo1, {
        "SUBTEMA_TITULO": "APIs REST: La Columna Vertebral de la Web Moderna",
        "CONTENIDO_TEXTO": (
            "API REST (Representational State Transfer) es un estilo de arquitectura para sistemas web:\n\n"
            "• Basada en recursos identificados por URLs (Ej: /productos, /usuarios/1).\n"
            "• Utiliza métodos HTTP estándar para interactuar con estos recursos.\n"
            "• Intercambia datos principalmente en formato JSON, aunque otros son posibles.\n"
            "• Es 'stateless' (sin estado): cada petición del cliente al servidor debe contener toda la información necesaria para entender la petición.\n\n"
            "Principios clave:\n"
            "  - Cliente-Servidor: Separación de responsabilidades.\n"
            "  - Sin estado (Stateless): Cada petición es independiente.\n"
            "  - Cacheable: Las respuestas pueden ser cacheadas.\n"
            "  - Sistema de capas: Puede haber intermediarios.\n"
            "  - Interfaz Uniforme: Uso consistente de URIs, métodos HTTP, etc."
        )
    })
    
    # Slide 2: Backend Simulado con JSON Server (Teoría)
    slide_desarrollo2 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo2, {
        "SUBTEMA_TITULO": "Simulando un Backend: JSON Server y su Utilidad",
        "CONTENIDO_TEXTO": (
            "JSON Server es una herramienta invaluable para el desarrollo frontend.\n\n"
            "• Permite crear rápidamente una API REST simulada a partir de un simple archivo JSON.\n"
            "• Ideal para prototipar y desarrollar sin depender de un backend real en etapas tempranas.\n"
            "• Facilita la prueba de la lógica de consumo de datos en Angular.\n\n"
            "Cómo funciona (conceptual):\n"
            "1. Se define la estructura de datos en un archivo `db.json` (ej: una lista de productos).\n"
            "2. Se ejecuta el comando `json-server --watch db.json` en la terminal.\n"
            "3. Automáticamente, se generan endpoints RESTful para cada recurso definido.\n"
            "   (Ej: `GET http://localhost:3000/productos` para obtener todos los productos).\n\n"
            "Esto nos permite enfocarnos en la lógica del frontend sin demoras y probar la interacción con una API de forma controlada."
        )
    })

    # Slide 3: Servicios en Angular: Concepto y Propósito (Teoría)
    slide_desarrollo3 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo3, {
        "SUBTEMA_TITULO": "Servicios en Angular: Centralizando la Lógica de Datos",
        "CONTENIDO_TEXTO": (
            "Los servicios son clases fundamentales en Angular para organizar y compartir código.\n\n"
            "Propósito principal:\n"
            "• Encapsular lógica de negocio o tareas específicas (ej: acceso a datos, autenticación, logging).\n"
            "• Compartir datos y funcionalidad entre componentes que no tienen una relación directa padre-hijo.\n"
            "• Promover un código más limpio, modular y reutilizable (Principio DRY - Don't Repeat Yourself).\n"
            "• Facilitar la inyección de dependencias, un patrón de diseño clave en Angular.\n\n"
            "En el contexto de APIs:\n"
            "• Son el lugar ideal para manejar todas las interacciones HTTP con el backend.\n"
            "• Utilizan el módulo `HttpClient` de Angular para realizar peticiones (GET, POST, PUT, DELETE).\n"
            "• Devuelven `Observables` (de la librería RxJS) para manejar respuestas asíncronas y flujos de datos."
        )
    })

    # Slide 4: Creando un Servicio Básico en Angular (Código Ligero)
    slide_desarrollo4 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo4, {
        "SUBTEMA_TITULO_CODIGO": "Estructura Esencial de un Servicio de Datos",
        "EJEMPLO_CODIGO": (
            "# 1. Generación con Angular CLI:\n"
            "ng generate service services/api\n\n"
            "# 2. Estructura básica (api.service.ts):\n"
            "import { Injectable } from '@angular/core';\n"
            "import { HttpClient } from '@angular/common/http';\n"
            "import { Observable } from 'rxjs';\n\n"
            "@Injectable({\n"
            "  providedIn: 'root' // Disponible en toda la aplicación\n"
            "})\n"
            "export class ApiService {\n"
            "  private apiUrl = 'http://localhost:3000/recurso'; // URL base de la API\n\n"
            "  constructor(private http: HttpClient) { }\n\n"
            "  // Ejemplo: Método para obtener datos (GET)\n"
            "  obtenerItems(): Observable<any[]> {\n"
            "    return this.http.get<any[]>(this.apiUrl);\n"
            "  }\n"
            "  // Otros métodos (POST, PUT, DELETE) seguirían un patrón similar.\n"
            "}"
        )
    })
    
    # Slide 5: Métodos HTTP y Operaciones CRUD (Teoría)
    slide_desarrollo5 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo5, {
        "SUBTEMA_TITULO": "Mapeo de Métodos HTTP a Operaciones CRUD",
        "CONTENIDO_TEXTO": (
            "Las APIs REST utilizan métodos HTTP para realizar operaciones sobre los recursos (datos).\n"
            "Estas operaciones se conocen comúnmente como CRUD (Create, Read, Update, Delete).\n\n"
            "• GET (Leer - Read):\n"
            "  - Solicita datos de un recurso específico o una colección.\n"
            "  - En Angular: `this.http.get('/productos')` o `this.http.get('/productos/1')`.\n"
            "  - Es una operación segura e idempotente (no cambia el estado del servidor).\n\n"
            "• POST (Crear - Create):\n"
            "  - Envía datos al servidor para crear un nuevo recurso.\n"
            "  - En Angular: `this.http.post('/productos', nuevoProducto)`.\n\n"
            "• PUT (Actualizar - Update):\n"
            "  - Reemplaza completamente un recurso existente con los datos enviados.\n"
            "  - En Angular: `this.http.put('/productos/1', datosActualizados)`.\n\n"
            "• DELETE (Eliminar - Delete):\n"
            "  - Elimina un recurso específico en el servidor.\n"
            "  - En Angular: `this.http.delete('/productos/1')`.\n\n"
            "(PATCH es otra opción para actualizaciones parciales, modificando solo algunos campos)."
        )
    })

    # Slide 6: Consumiendo Servicios: El Rol de los Componentes (Teoría)
    slide_desarrollo6 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo6, {
        "SUBTEMA_TITULO": "Integración: Componentes Consumiendo Servicios",
        "CONTENIDO_TEXTO": (
            "Los componentes son los responsables de presentar los datos y manejar la interacción del usuario.\n\n"
            "Flujo típico de consumo de un servicio en un componente:\n"
            "1. Inyección del Servicio: El componente declara el servicio como una dependencia en su constructor. Angular se encarga de proveer la instancia.\n"
            "   `constructor(private miServicio: ApiService) { }`\n\n"
            "2. Llamada al Método del Servicio: Generalmente en el hook `ngOnInit` (para carga inicial) o en respuesta a un evento del usuario (ej: clic en un botón).\n"
            "   `this.miServicio.obtenerItems()`\n\n"
            "3. Suscripción al Observable: Los métodos de `HttpClient` devuelven Observables. El componente se 'suscribe' para recibir los datos cuando estén disponibles (o un error si la petición falla).\n"
            "   `.subscribe(\n     datos => { this.items = datos; /* Lógica con los datos */ },\n     error => { console.error('Error:', error); /* Manejo de error */ }\n   );`\n\n"
            "4. Actualización de la Vista: Los datos recibidos se asignan a propiedades del componente, que luego se muestran en la plantilla HTML usando data binding."
        )
    })

    # Slide 7: Herramientas para Probar y Depurar APIs (Teoría)
    slide_desarrollo7 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo7, {
        "SUBTEMA_TITULO": "Herramientas Esenciales para Probar APIs",
        "CONTENIDO_TEXTO": (
            "Probar las APIs es crucial antes y durante la integración con el frontend.\n\n"
            "Herramientas populares:\n"
            "• Postman: Una completa plataforma gráfica para diseñar, construir, probar y documentar APIs. Permite enviar todo tipo de peticiones HTTP, organizar colecciones y automatizar pruebas.\n\n"
            "• Insomnia: Similar a Postman, con un enfoque en la simplicidad y la velocidad. Buena alternativa para pruebas rápidas y diseño de APIs.\n\n"
            "• Thunder Client (Extensión VS Code): Permite probar APIs directamente desde el editor de código, muy conveniente para desarrolladores que prefieren no salir de su entorno.\n\n"
            "• Herramientas de Desarrollador del Navegador (Pestaña Network): Útil para inspeccionar las peticiones HTTP que realiza tu aplicación Angular y las respuestas del servidor en tiempo real."
        )
    })
    
    # --- Momento 4: Práctica ---
    slide_practica = crear_diapositiva_con_marcadores(prs, "Layout_Practica")
    reemplazar_marcadores(slide_practica, {
        "TITULO_PRACTICA": "Práctica: App de Contactos con API REST (JSON Server)",
        "INSTRUCCIONES_PRACTICA": (
             "Actividad Práctica Guiada (Conceptual y Esencial):\n\n"
             "Objetivo: Implementar las operaciones básicas para gestionar una lista de contactos consumiendo una API simulada con JSON Server.\n\n"
             "Pasos Clave:\n"
             "1. Configurar JSON Server:\n"
             "   - Crear un archivo `db.json` con una lista inicial de contactos (id, nombre, email, telefono).\n"
             "   - Iniciar el servidor: `json-server --watch db.json`.\n\n"
             "2. En el proyecto Angular:\n"
             "   - Asegurar que `HttpClientModule` esté importado en `app.module.ts`.\n"
             "   - Generar un servicio (`ContactService`).\n"
             "   - Implementar en `ContactService` métodos para:\n"
             "     - Obtener todos los contactos (GET).\n"
             "     - Crear un nuevo contacto (POST).\n"
             "     - (Opcional) Actualizar y Eliminar.\n"
             "   - Crear un componente (`ContactListComponent`) para:\n"
             "     - Inyectar y usar `ContactService`.\n"
             "     - Mostrar la lista de contactos.\n"
             "     - Tener un formulario simple para añadir un nuevo contacto."
        ),
        "RECURSOS_NECESARIOS": (
             "• Proyecto Angular existente.\n"
             "• JSON Server instalado (`npm install -g json-server`).\n"
             "• Editor de código.\n"
             "• Navegador web y sus herramientas de desarrollo.\n"
             "• Tiempo estimado: 30-40 minutos."
        )
    })
    
    # --- Momento 5: Cierre ---
    slide_cierre = crear_diapositiva_con_marcadores(prs, "Layout_Cierre")
    reemplazar_marcadores(slide_cierre, {
        "TITULO_CIERRE": "Resumen: Conectando Angular con APIs REST",
        "RESUMEN_IDEAS": (
             "En esta sesión hemos explorado:\n\n"
             "• Los conceptos fundamentales de APIs REST y su arquitectura.\n"
             "• Cómo simular un backend rápidamente con JSON Server para desarrollo.\n"
             "• La importancia y estructura de los Servicios en Angular para encapsular lógica de datos.\n"
             "• La implementación de operaciones CRUD usando `HttpClient` y métodos HTTP.\n"
             "• El flujo de consumo de servicios desde componentes mediante Observables.\n"
             "• Herramientas útiles para probar y depurar APIs."
        ),
        "PREGUNTAS_METACOGNITIVAS": (
             "Reflexionemos:\n\n"
             "1. ¿Qué ventajas principales ofrece separar la lógica de acceso a datos en servicios dedicados?\n"
             "2. ¿Cómo facilita JSON Server el desarrollo de aplicaciones que dependen de un backend?\n"
             "3. ¿Por qué es crucial entender el concepto de operaciones asíncronas (Observables) al trabajar con APIs?\n"
             "4. ¿Qué desafíos podrías anticipar al conectar tu aplicación Angular con APIs reales en un entorno de producción?"
        )
    })
    
    # Diapositiva de Preguntas
    slide_preguntas = crear_diapositiva_con_marcadores(prs, "Layout_Preguntas")
    reemplazar_marcadores(slide_preguntas, {
        "TITULO_PREGUNTAS": "¿Dudas sobre Servicios, APIs REST o JSON Server?",
        "DUDAS_LISTADO": (
             "Resolvamos dudas sobre:\n\n"
             "• Conceptos de REST, recursos y métodos HTTP.\n"
             "• Configuración y uso de JSON Server.\n"
             "• Creación, inyección y propósito de los servicios en Angular.\n"
             "• El uso de `HttpClient` y el manejo de `Observables` y suscripciones.\n"
             "• Estrategias básicas para el tratamiento de errores en peticiones HTTP.\n"
             "• Buenas prácticas para estructurar la comunicación con APIs.\n\n"
             "¡Aclaremos estos conceptos para poder avanzar con confianza!"
        )
    })
    
    # Diapositiva Final
    slide_final = crear_diapositiva_con_marcadores(prs, "1_Layout_Preguntas")
    reemplazar_marcadores(slide_final, {
        "TITULO_AGRADECIMIENTO": "¡Felicidades por dar el paso hacia aplicaciones web conectadas!\n\nEn la próxima sesión, integraremos todo lo aprendido para crear una aplicación web completa con múltiples funcionalidades."
    })

    # Guardar presentación
    os.makedirs(RUTA_PRESENTACIONES, exist_ok=True)

    nombre_archivo = 'presentacion_S09_s1_servicios_angular_rest_a.pptx' # Manteniendo el nombre original o usar _v2 si se prefiere
    ruta_salida = os.path.join(RUTA_PRESENTACIONES, nombre_archivo)
    try:
        prs.save(ruta_salida)
        print(f"Presentación (v2 Teórico) de la Semana 9 - Sesión 1 generada exitosamente en: {ruta_salida}")
    except Exception as e:
        print(f"Error al guardar la presentación (v2 Teórico): {e}")
   

def generar_presentacion_semana11_1_javascript():
    """
    Genera una presentación completa para la semana 11, sesión 1 del curso de 
    JavaScript Avanzado.
    Tema: Introducción a React: concepto, instalación, creación de proyectos, JSX, TSX.
    """
    print("Iniciando generación de presentación Semana 11 - Sesión 1 (JavaScript Avanzado - React Intro)")
    
    # Cargar la plantilla y los datos del curso
    ruta_plantilla = os.path.join(RUTA_PLANTILLAS, 'plantilla_2.pptx')
    ruta_json = os.path.join(RUTA_JSON, 'JavaScript.json')

    if not os.path.exists(ruta_plantilla):
        print(f"Error: No se encontró la plantilla en {ruta_plantilla}")
        return
    
    if not os.path.exists(ruta_json):
        print(f"Error: No se encontró el archivo JSON en {ruta_json}")
        return

    prs = Presentation(ruta_plantilla)
    datos_curso_completo = cargar_json(ruta_json) 
    if not datos_curso_completo or 'silabo' not in datos_curso_completo: 
        print("Error: No se pudieron cargar los datos del curso o estructura JSON inválida.")
        return
    datos_curso = datos_curso_completo['silabo'] # Renamed for consistency
    docente_nombre = "Ms. Johan Max A. López Heredia" # Direct assignment

    # --- Diapositiva de Portada ---
    slide_portada = crear_diapositiva_con_marcadores(prs, "Layout_Portada_UTP_1")
    titulo_portada = "S11.s1 - Introducción a React: Conceptos y Primeros Pasos" 
    reemplazar_marcadores(slide_portada, {
        "TITULO_SEMANA": titulo_portada,
        "NOMBRE_SEMANA": f"{datos_curso.get('curso', 'Curso Desconocido')} - {datos_curso.get('codigo', 'Código Desconocido')}",
        "DOCENTE_NOMBRE": docente_nombre
    })
    
    # --- Momento 1: Inicio ---
    slide_inicio = crear_diapositiva_con_marcadores(prs, "Layout_Inicio")
    reemplazar_marcadores(slide_inicio, {
        "TITULO_INICIO": "Construyendo Interfaces Dinámicas con React",
        "ACTIVIDAD_ROMPEHIELO": (
            "Actividad inicial (5 minutos):\n\n"
            "Piensa en cómo usas Instagram o Facebook:\n\n"
            "- Cuando das 'Me gusta' a una foto, ¿se recarga toda la página o solo cambia el corazón?\n"
            "- Cuando escribes un comentario, ¿se actualiza solo esa parte?\n"
            "- ¿Cómo crees que logran que la página sea tan interactiva y rápida?\n\n"
            "React es una herramienta que ayuda a construir este tipo de interfaces 'reactivas'."
        )
    })
    
    # --- Momento 2: Utilidad ---
    # Diapositiva de Logro
    slide_logro = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_Logro")
    logro_sesion = ("Al finalizar la sesión, el estudiante comprende los conceptos fundamentales de React, "
                    "configura un entorno de desarrollo básico, crea su primer proyecto y reconoce la sintaxis "
                    "JSX/TSX para la construcción de componentes de interfaz.")
    
    reemplazar_marcadores(slide_logro, {
        "ENCABEZADO_LOGRO": "Logro de Aprendizaje de la Sesión",
        "LOGRO_TEXTO": logro_sesion
    })
    
    # Diapositiva de Utilidad y Dudas
    slide_dudas = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_dudas")
    reemplazar_marcadores(slide_dudas, {
        "ENCABEZADO_DUDAS": "Dudas Frecuentes / Repaso",
        "DUDAS_TEXTO": (
            "- ¿React es otro framework como Angular? ¿Cuál es la diferencia?\n"
            "- ¿Es solo para interfaces o hace más cosas?\n"
            "- ¿Qué es JSX/TSX? ¿Es HTML o JavaScript?\n"
            "- ¿Necesito aprenderlo si ya sé JavaScript básico?\n"
            "- ¿Facebook usa React?"
        ),
        "ENCABEZADO_IMPORTANCIA_TEMA": "Importancia de Aprender React",
        "IMPORTANCIA_TEXTO": (
            "React es una habilidad muy valiosa en el desarrollo web actual porque:\n\n"
            "• Es una de las librerías/frameworks de JavaScript más populares y demandadas.\n"
            "• Permite crear interfaces de usuario (UI) interactivas y reutilizables (componentes).\n"
            "• Facilita la creación de Single Page Applications (SPAs) eficientes.\n"
            "• Tiene un gran ecosistema de herramientas y librerías complementarias.\n"
            "• Es mantenido por Meta (Facebook) y una gran comunidad de desarrolladores."
        )
    })
    
    # Diapositiva de Saberes Previos
    slide_saberes = crear_diapositiva_con_marcadores(prs, "Layout_Saberes")
    reemplazar_marcadores(slide_saberes, {
        "TITULO_CONOCIMIENTOS": "Recordando el Frontend",
        "DINAMICA_PREVIA": (
            "Recordemos brevemente (5 minutos):\n\n"
            "1. ¿Qué tres tecnologías principales forman la base del frontend web? (HTML, CSS, JavaScript).\n\n"
            "2. ¿Para qué sirve cada una de ellas?\n"
            "   - HTML: Estructura y contenido.\n"
            "   - CSS: Estilo y presentación visual.\n"
            "   - JavaScript: Comportamiento e interactividad.\n\n"
            "3. ¿Qué es el DOM (Document Object Model)? (La representación del HTML que JavaScript puede manipular).\n\n"
            "React nos ayuda a manipular el DOM de forma más eficiente y declarativa."
        )
    })
    
    # --- Momento 3: Transformación ---
    # Diapositiva 1: ¿Qué es React?
    slide_desarrollo1 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo1, {
        "SUBTEMA_TITULO": "¿Qué es React?",
        "CONTENIDO_TEXTO": (
            "React es una librería de JavaScript (aunque a menudo se le considera un framework) para construir interfaces de usuario (UI).\n\n"
            "• Se enfoca en la 'V' (Vista) del patrón Modelo-Vista-Controlador (MVC).\n"
            "• Su idea central es construir UIs a partir de piezas reutilizables llamadas Componentes.\n"
            "• Utiliza un DOM Virtual para optimizar las actualizaciones de la pantalla, haciéndolo muy eficiente.\n"
            "• Es declarativo: Le dices a React CÓMO quieres que se vea la UI con ciertos datos, y React se encarga de actualizarla cuando los datos cambian.\n\n"
            "Fue creado por Facebook y es de código abierto."
        )
    })
    
    # Diapositiva 2: Instalación y Creación de Proyectos (Básico)
    slide_desarrollo2 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo2, {
        "SUBTEMA_TITULO_CODIGO": "Preparando el Entorno React (Básico)",
        # Comandos básicos, sin entrar en detalles de npx vs npm global.
        "EJEMPLO_CODIGO": (
            "# Se necesita NodeJS y npm (o yarn) instalados previamente.\n\n"
            "# Método recomendado para crear una nueva app React:\n"
            "npx create-react-app mi-app-react\n"
            "# (npx ejecuta el paquete sin instalarlo globalmente).\n\n"
            "# Esto crea una carpeta 'mi-app-react' con toda la estructura inicial.\n\n"
            "# Entrar a la carpeta del proyecto:\n"
            "cd mi-app-react\n\n"
            "# Ejecutar la aplicación en modo desarrollo:\n"
            "npm start\n"
            "# (Esto abrirá la app en http://localhost:3000 en tu navegador).\n\n"
            "# Alternativa moderna (más rápida): Vite\n"
            "# npm create vite@latest mi-app-vite -- --template react-ts"
        )
    })

    # Diapositiva 3: ¿Qué es JSX y TSX?
    slide_desarrollo3 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo3, {
        "SUBTEMA_TITULO": "JSX / TSX: HTML dentro de JavaScript/TypeScript",
        "CONTENIDO_TEXTO": (
            "JSX (JavaScript XML) y TSX (TypeScript XML) son extensiones de sintaxis que permiten escribir 'HTML' directamente dentro de tu código JavaScript o TypeScript.\n\n"
            "• NO es HTML real en el navegador, pero se parece mucho.\n"
            "• Se transpila (convierte) a llamadas de funciones JavaScript normales (`React.createElement()`) durante el proceso de construcción.\n\n"
            "¿Por qué usarlo?\n"
            "- Hace que el código de los componentes sea más legible y fácil de entender cómo se verá la UI.\n"
            "- Permite combinar la lógica de JavaScript (variables, funciones) directamente con la estructura de la UI.\n\n"
            "Es una característica central de React."
        )
    })

    # Diapositiva 4: Ejemplo Básico de JSX/TSX
    slide_desarrollo4 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo4, {
        "SUBTEMA_TITULO_CODIGO": "Ejemplo Simple de JSX/TSX",
        # Código conceptual JSX/TSX.
        "EJEMPLO_CODIGO": (
            "// En un archivo de componente React (ej. App.tsx o App.js)\n\n"
            "function App() {\n"
            "  const nombreUsuario = \"Ana\";\n"
            "  const titulo = <h1>¡Bienvenida, {nombreUsuario}!</h1>;\n\n"
            "  return (\n"
            "    <div>\n"
            "      {titulo}\n"
            "      <p>Esta es tu primera app con React.</p>\n"
            "    </div>\n"
            "  );\n"
            "}\n\n"
            "// export default App;\n\n"
            "# Notas:\n"
            "# - Las etiquetas HTML se escriben directamente.\n"
            "# - Las expresiones JavaScript (variables, funciones) van entre llaves `{}`.\n"
            "# - `className` se usa en lugar de `class` para atributos CSS."
        )
    })

    # Diapositiva 5: Componentes: Los Bloques de React
    slide_desarrollo5 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo5, {
        "SUBTEMA_TITULO": "Componentes: Los Ladrillos de tu UI",
        "CONTENIDO_TEXTO": (
            "La idea principal de React es construir UIs complejas a partir de piezas pequeñas, independientes y reutilizables llamadas Componentes.\n\n"
            "• Un componente es como una función de JavaScript que devuelve elementos React (JSX/TSX) para describir qué debe aparecer en la pantalla.\n\n"
            "Tipos de Componentes (principalmente):\n"
            "• Componentes Funcionales: Son funciones JavaScript/TypeScript. Son la forma moderna y recomendada de escribir componentes.\n"
            "• Componentes de Clase (más antiguos, aún existentes).\n\n"
            "Ejemplos: Un botón, un campo de formulario, una tarjeta de producto, una barra de navegación, ¡incluso toda una página puede ser un componente!"
        )
    })
    
    # --- Momento 4: Práctica ---
    slide_practica = crear_diapositiva_con_marcadores(prs, "Layout_Practica")
    reemplazar_marcadores(slide_practica, {
        "TITULO_PRACTICA": "Práctica: Creando Nuestro Primer Proyecto React",
        "INSTRUCCIONES_PRACTICA": (
             "Actividad Práctica Guiada (Sencilla):\n\n"
             "Objetivo: Crear un proyecto React básico y ver la estructura inicial.\n\n"
             "Pasos:\n"
             "1. Asegúrate de tener NodeJS y npm (o yarn) instalados.\n"
             "2. Abre una terminal o línea de comandos.\n"
             "3. Navega a una carpeta donde quieras crear tu proyecto.\n"
             "4. Ejecuta el comando para crear una nueva aplicación React:\n"
             "   `npx create-react-app mi-primera-app-react`\n"
             "   (Esto tomará unos minutos mientras descarga e instala todo).\n\n"
             "5. Una vez terminado, entra a la carpeta del proyecto:\n"
             "   `cd mi-primera-app-react`\n\n"
             "6. Inicia el servidor de desarrollo:\n"
             "   `npm start`\n\n"
             "7. Se abrirá una pestaña en tu navegador (usualmente http://localhost:3000) mostrando la página de inicio de React.\n\n"
             "8. (Opcional) Abre la carpeta `mi-primera-app-react` en VS Code y explora el archivo `src/App.js` (o `App.tsx`). ¡Intenta cambiar el texto que se muestra y guarda para ver los cambios en el navegador!\n\n"
             "¡El objetivo es crear y ejecutar la aplicación base de React!"
        ),
        "RECURSOS_NECESARIOS": (
             "• NodeJS y npm (o yarn) instalados.\n"
             "• Terminal o línea de comandos.\n"
             "• Navegador web.\n"
             "• (Opcional) Editor de código como VS Code.\n"
             "• Tiempo estimado: 20-30 minutos (la creación del proyecto puede tardar)."
        )
    })
    
    # --- Momento 5: Cierre ---
    # Diapositiva de Cierre y Síntesis
    slide_cierre = crear_diapositiva_con_marcadores(prs, "Layout_Cierre")
    reemplazar_marcadores(slide_cierre, {
        "TITULO_CIERRE": "Resumen: Primeros Pasos con React",
        "RESUMEN_IDEAS": (
             "En esta sesión hemos:\n\n"
             "• Introducido React como una librería para construir interfaces de usuario.\n"
             "• Visto cómo crear un nuevo proyecto React usando `create-react-app`.\n"
             "• Comprendido qué es JSX/TSX y cómo permite escribir HTML-como sintaxis en JavaScript/TypeScript.\n"
             "• Presentado el concepto fundamental de Componentes como bloques de construcción de UIs.\n"
             "• Realizado una práctica para crear y ejecutar nuestra primera aplicación React."
        ),
        "PREGUNTAS_METACOGNITIVAS": (
             "Reflexionemos:\n\n"
             "1. ¿Qué te parece la idea de escribir 'HTML' dentro de JavaScript (JSX/TSX)?\n"
             "2. ¿Qué ventajas crees que tiene construir interfaces a partir de componentes reutilizables?\n"
             "3. Comparado con JavaScript puro o jQuery (si lo conoces), ¿qué te llama la atención de React hasta ahora?"
        )
    })
    
    # Diapositiva de Preguntas
    slide_preguntas = crear_diapositiva_con_marcadores(prs, "Layout_Preguntas")
    reemplazar_marcadores(slide_preguntas, {
        "TITULO_PREGUNTAS": "¿Dudas sobre Conceptos Básicos de React?",
        "DUDAS_LISTADO": (
             "Ahora es el momento de resolver dudas sobre:\n\n"
             "• Qué es React y su propósito principal.\n"
             "• El proceso de instalación y creación de proyectos (`create-react-app`).\n"
             "• La sintaxis de JSX/TSX y cómo se usa.\n"
             "• El concepto inicial de Componentes.\n"
             "• Cualquier problema encontrado durante la práctica.\n\n"
             "¡Asegurémonos de tener claros estos fundamentos!"
        )
    })
    
    # Diapositiva Final de Agradecimiento
    slide_final = crear_diapositiva_con_marcadores(prs, "1_Layout_Preguntas")
    reemplazar_marcadores(slide_final, {
        "TITULO_AGRADECIMIENTO": "¡Felicidades por iniciar tu viaje con React!\n\nEn la próxima sesión, profundizaremos en los Componentes de React, cómo pasarles datos (props) y manejar su estado interno."
    })

    # --- FIN DE LAS DIAPOSITIVAS ---

    # Verificar contenido final (opcional pero útil para depurar)
    # verificar_contenido_diapositivas(prs)

    # Crear directorio si no existe
    os.makedirs(RUTA_PRESENTACIONES, exist_ok=True)

    # Guardar la presentación generada
    nombre_archivo = 'presentacion_S11_s1_introduccion_react_jsx_tsx.pptx'
    ruta_salida = os.path.join(RUTA_PRESENTACIONES, nombre_archivo)
    try:
        prs.save(ruta_salida)
        print(f"Presentación de la Semana 11 - Sesión 1 (JavaScript Avanzado) generada exitosamente en: {ruta_salida}")
    except Exception as e:
        print(f"Error al guardar la presentación: {e}")

# --- Fin de la función ---


def generar_presentacion_semana12_1_javascript():
    """
    Genera una presentación completa para la semana 12, sesión 1 del curso de 
    JavaScript Avanzado.
    Tema: Elementos de React: Componentes, Navegación y Formularios.
    """
    print("Iniciando generación de presentación Semana 12 - Sesión 1 (JavaScript Avanzado - Elementos de React)")
    
    # Cargar la plantilla y los datos del curso
    ruta_plantilla = os.path.join(RUTA_PLANTILLAS, 'plantilla_2.pptx')
    ruta_json = os.path.join(RUTA_JSON, 'JavaScript.json')

    if not os.path.exists(ruta_plantilla):
        print(f"Error: No se encontró la plantilla en {ruta_plantilla}")
        return
    
    if not os.path.exists(ruta_json):
        print(f"Error: No se encontró el archivo JSON en {ruta_json}")
        return

    prs = Presentation(ruta_plantilla)
    datos_curso_completo = cargar_json(ruta_json) 
    if not datos_curso_completo or 'silabo' not in datos_curso_completo: 
        print("Error: No se pudieron cargar los datos del curso o estructura JSON inválida.")
        return
    datos_curso = datos_curso_completo['silabo']
    docente_nombre = "Ms. Johan Max A. López Heredia"

    # --- Diapositiva de Portada ---
    slide_portada = crear_diapositiva_con_marcadores(prs, "Layout_Portada_UTP_1")
    titulo_portada = "S12.s1 - Elementos de React: Componentes, Navegación y Formularios" 
    reemplazar_marcadores(slide_portada, {
        "TITULO_SEMANA": titulo_portada,
        "NOMBRE_SEMANA": f"{datos_curso.get('curso', 'Curso Desconocido')} - {datos_curso.get('codigo', 'Código Desconocido')}",
        "DOCENTE_NOMBRE": docente_nombre
    })
    
    # --- Momento 1: Inicio ---
    slide_inicio = crear_diapositiva_con_marcadores(prs, "Layout_Inicio")
    reemplazar_marcadores(slide_inicio, {
        "TITULO_INICIO": "Construyendo los Cimientos de una Aplicación React",
        "ACTIVIDAD_ROMPEHIELO": (
            "Actividad inicial (5 minutos):\n\n"
            "Piensa en una aplicación web compleja que uses (ej. un e-commerce, una red social):\n\n"
            "- ¿Cómo pasas de una sección a otra (de 'productos' al 'carrito')?\n"
            "- ¿Qué partes de la página se repiten (botones, barras de navegación)?\n"
            "- ¿Cómo llenas la información para registrarte o comprar?\n\n"
            "Hoy aprenderemos a crear esos 'bloques de construcción' (componentes), a navegar entre ellos y a manejar formularios con React."
        )
    })
    
    # --- Momento 2: Utilidad ---
    slide_logro = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_Logro")
    logro_sesion = ("Al finalizar la sesión, el estudiante implementa componentes funcionales, "
                    "formularios controlados y un sistema de navegación básico en una aplicación web "
                    "utilizando React, para estructurar interfaces de usuario interactivas.")
    
    reemplazar_marcadores(slide_logro, {
        "ENCABEZADO_LOGRO": "Logro de Aprendizaje de la Sesión",
        "LOGRO_TEXTO": logro_sesion
    })
    
    slide_dudas = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_dudas")
    reemplazar_marcadores(slide_dudas, {
        "ENCABEZADO_DUDAS": "Dudas Frecuentes / Repaso",
        "DUDAS_TEXTO": (
            "- ¿Cómo se comunican los componentes entre sí?\n"
            "- ¿Qué es el 'estado' de un componente?\n"
            "- ¿Cómo hago para que mi app tenga diferentes 'páginas' si es una SPA?\n"
            "- ¿Es muy diferente manejar formularios en React que en HTML puro?\n"
            "- ¿Qué es React Router?"
        ),
        "ENCABEZADO_IMPORTANCIA_TEMA": "Importancia de los Elementos de React",
        "IMPORTANCIA_TEXTO": (
            "Dominar componentes, navegación y formularios es esencial porque:\n\n"
            "• Son los tres pilares sobre los que se construye CUALQUIER aplicación React.\n"
            "• Permiten crear UIs complejas, interactivas y escalables.\n"
            "• La navegación del lado del cliente (React Router) define la experiencia de usuario en una SPA.\n"
            "• El manejo de formularios es clave para cualquier aplicación que requiera entrada de datos del usuario.\n"
            "• La componentización es la base de la reutilización y mantenibilidad del código en React."
        )
    })
    
    slide_saberes = crear_diapositiva_con_marcadores(prs, "Layout_Saberes")
    reemplazar_marcadores(slide_saberes, {
        "TITULO_CONOCIMIENTOS": "Recordando Nuestra Introducción a React",
        "DINAMICA_PREVIA": (
            "Recordemos brevemente la sesión anterior (5 minutos):\n\n"
            "1. ¿Qué es React? (Una librería para construir interfaces de usuario).\n\n"
            "2. ¿Cuál es la idea principal de React? (Construir UIs con piezas reutilizables llamadas Componentes).\n\n"
            "3. ¿Qué es JSX? (Una sintaxis que nos permite escribir 'HTML' en JavaScript).\n\n"
            "4. ¿Qué comando usamos para crear una nueva app de React? (`npx create-react-app` o `npm create vite`).\n\n"
            "Hoy profundizaremos en cómo hacer que esos componentes sean dinámicos, cómo navegar entre ellos y cómo capturar datos del usuario."
        )
    })
    
    # --- Momento 3: Transformación ---
    slide_desarrollo1 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo1, {
        "SUBTEMA_TITULO": "Componentes en React: Props y Estado",
        "CONTENIDO_TEXTO": (
            "Los componentes son el corazón de React. Tienen dos conceptos clave:\n\n"
            "1. Props (Propiedades):\n"
            "   - Son la forma de pasar datos de un componente padre a un componente hijo.\n"
            "   - Son de solo lectura (inmutables) dentro del componente hijo.\n"
            "   - Ejemplo: `<PerfilUsuario nombre=\"Ana\" edad={30} />`\n\n"
            "2. State (Estado):\n"
            "   - Es información que un componente mantiene y puede cambiar con el tiempo.\n"
            "   - Cuando el estado cambia, React vuelve a renderizar el componente.\n"
            "   - Se maneja con el Hook `useState`.\n"
            "   - Ejemplo: Un contador, el texto de un input, si un menú está abierto o cerrado."
        )
    })
    
    slide_desarrollo2 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo2, {
        "SUBTEMA_TITULO_CODIGO": "El Hook `useState`: Dando Memoria a los Componentes",
        "EJEMPLO_CODIGO": (
            "// Importamos useState desde React\n"
            "import { useState } from 'react';\n\n"
            "function Contador() {\n"
            "  // Declaramos una variable de estado llamada 'conteo'\n"
            "  // y una función para actualizarla, 'setConteo'.\n"
            "  const [conteo, setConteo] = useState(0); // 0 es el valor inicial\n\n"
            "  return (\n"
            "    <div>\n"
            "      <p>Has hecho clic {conteo} veces</p>\n"
            "      <button onClick={() => setConteo(conteo + 1)}>\n"
            "        Hazme clic\n"
            "      </button>\n"
            "    </div>\n"
            "  );\n"
            "}"
        )
    })

    slide_desarrollo3 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo3, {
        "SUBTEMA_TITULO": "Manejo de Formularios: Formularios Controlados",
        "CONTENIDO_TEXTO": (
            "En React, la técnica común para manejar formularios se llama 'componentes controlados'.\n\n"
            "La idea es que el estado del componente React es la 'única fuente de la verdad'.\n\n"
            "Pasos:\n"
            "1. Crear una variable de estado en el componente para cada input del formulario.\n"
            "   `const [nombre, setNombre] = useState('');`\n\n"
            "2. Vincular el valor del input al estado.\n"
            "   `<input type=\"text\" value={nombre} />`\n\n"
            "3. Crear una función que actualice el estado cada vez que el input cambia (evento `onChange`).\n"
            "   `<input ... onChange={(e) => setNombre(e.target.value)} />`\n\n"
            "Así, React siempre 'sabe' cuál es el valor del formulario."
        )
    })

    slide_desarrollo4 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo4, {
        "SUBTEMA_TITULO_CODIGO": "Ejemplo de Formulario Controlado",
        "EJEMPLO_CODIGO": (
            "function FormularioNombre() {\n"
            "  const [nombre, setNombre] = useState('');\n\n"
            "  const handleSubmit = (event) => {\n"
            "    event.preventDefault(); // Evita que la página se recargue\n"
            "    alert(`Un nombre fue enviado: ${nombre}`);\n"
            "  }\n\n"
            "  return (\n"
            "    <form onSubmit={handleSubmit}>\n"
            "      <label>Nombre:</label>\n"
            "      <input \n"
            "        type=\"text\" \n"
            "        value={nombre} \n"
            "        onChange={(e) => setNombre(e.target.value)} \n"
            "      />\n"
            "      <button type=\"submit\">Enviar</button>\n"
            "    </form>\n"
            "  );\n"
            "}"
        )
    })

    slide_desarrollo5 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo5, {
        "SUBTEMA_TITULO": "Navegación en React con React Router",
        "CONTENIDO_TEXTO": (
            "React no incluye un sistema de rutas por defecto. La solución estándar de la comunidad es la librería `React Router`.\n\n"
            "Permite crear una experiencia de Single Page Application (SPA) donde la navegación entre 'páginas' ocurre sin recargar el navegador.\n\n"
            "Pasos para instalar:\n"
            "`npm install react-router-dom`\n\n"
            "Componentes principales que usaremos:\n"
            "• `BrowserRouter`: Envuelve toda la aplicación para habilitar el enrutamiento.\n"
            "• `Routes`: Contenedor para todas las rutas individuales.\n"
            "• `Route`: Define una ruta, asociando una URL (path) a un componente.\n"
            "• `Link`: Reemplaza a la etiqueta `<a>` para crear enlaces de navegación interna."
        )
    })

    slide_desarrollo6 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo6, {
        "SUBTEMA_TITULO_CODIGO": "Configurando Rutas Básicas",
        "EJEMPLO_CODIGO": (
            "// En tu archivo principal (index.js o App.js)\n"
            "import { BrowserRouter, Routes, Route } from 'react-router-dom';\n"
            "import HomePage from './HomePage';\n"
            "import AboutPage from './AboutPage';\n\n"
            "function App() {\n"
            "  return (\n"
            "    <BrowserRouter>\n"
            "      <Routes>\n"
            "        <Route path=\"/\" element={<HomePage />} />\n"
            "        <Route path=\"/about\" element={<AboutPage />} />\n"
            "      </Routes>\n"
            "    </BrowserRouter>\n"
            "  );\n"
            "}\n\n"
            "// En otro componente (ej. una barra de navegación)\n"
            "import { Link } from 'react-router-dom';\n\n"
            "<nav>\n"
            "  <Link to=\"/\">Inicio</Link>\n"
            "  <Link to=\"/about\">Acerca de</Link>\n"
            "</nav>"
        )
    })
    
    # --- Momento 4: Práctica ---
    slide_practica = crear_diapositiva_con_marcadores(prs, "Layout_Practica")
    reemplazar_marcadores(slide_practica, {
        "TITULO_PRACTICA": "Práctica: Creando una Mini-App con Navegación y Formulario",
        "INSTRUCCIONES_PRACTICA": (
             "Actividad Práctica Guiada:\n\n"
             "Objetivo: Crear una aplicación de dos páginas con un formulario simple.\n\n"
             "Pasos:\n"
             "1. En tu proyecto React, instala React Router: `npm install react-router-dom`.\n"
             "2. Crea dos nuevos componentes: `HomePage.js` y `ContactPage.js`.\n"
             "   - `HomePage` solo mostrará un título 'Bienvenido a la Página de Inicio'.\n"
             "3. En `App.js`, configura las rutas para que `/` muestre `HomePage` y `/contact` muestre `ContactPage`.\n"
             "4. Crea un componente `Navbar.js` con enlaces (`<Link>`) a las dos rutas.\n"
             "5. Incluye el componente `<Navbar />` en `App.js` para que sea visible en todas las páginas.\n"
             "6. En `ContactPage.js`, crea un formulario controlado con un campo para 'email'.\n"
             "7. El formulario debe tener un estado para guardar el email y un botón de envío que muestre una alerta con el email ingresado.\n\n"
             "¡El objetivo es tener una app donde puedas navegar entre dos páginas y usar un formulario funcional!"
        ),
        "RECURSOS_NECESARIOS": (
             "• Proyecto React creado (con `create-react-app` o `vite`).\n"
             "• Terminal para instalar paquetes.\n"
             "• Editor de código.\n"
             "• Tiempo estimado: 30-40 minutos."
        )
    })
    
    # --- Momento 5: Cierre ---
    slide_cierre = crear_diapositiva_con_marcadores(prs, "Layout_Cierre")
    reemplazar_marcadores(slide_cierre, {
        "TITULO_CIERRE": "Resumen: Los Pilares de una App React",
        "RESUMEN_IDEAS": (
             "En esta sesión hemos aprendido a:\n\n"
             "• Dar 'memoria' a nuestros componentes con el estado (`useState`).\n"
             "• Pasar datos a componentes usando `props`.\n"
             "• Capturar la entrada del usuario con formularios controlados.\n"
             "• Estructurar una aplicación de múltiples páginas usando React Router.\n"
             "• Crear enlaces de navegación que no recargan la página.\n\n"
             "Estos son los conceptos fundamentales para construir cualquier aplicación web con React."
        ),
        "PREGUNTAS_METACOGNITIVAS": (
             "Reflexionemos:\n\n"
             "1. ¿Por qué es importante que el estado sea la 'fuente de la verdad' en un formulario controlado?\n"
             "2. ¿Qué ventajas le ves a la navegación de React Router comparada con la navegación tradicional de links `<a>`?\n"
             "3. ¿En qué otro tipo de componente (además de un formulario) crees que sería útil usar el estado (`useState`)?"
        )
    })
    
    slide_preguntas = crear_diapositiva_con_marcadores(prs, "Layout_Preguntas")
    reemplazar_marcadores(slide_preguntas, {
        "TITULO_PREGUNTAS": "¿Dudas sobre Componentes, Estado, Formularios o Rutas?",
        "DUDAS_LISTADO": (
             "Ahora es el momento de resolver dudas sobre:\n\n"
             "• La diferencia entre `props` y `state`.\n"
             "• El funcionamiento del hook `useState`.\n"
             "• Cómo manejar diferentes tipos de inputs en formularios.\n"
             "• La configuración de `React Router`.\n"
             "• Cualquier problema encontrado durante la práctica.\n\n"
             "¡Asegurémonos de que estos cimientos estén sólidos!"
        )
    })
    
    slide_final = crear_diapositiva_con_marcadores(prs, "1_Layout_Preguntas")
    reemplazar_marcadores(slide_final, {
        "TITULO_AGRADECIMIENTO": "¡Felicidades por construir una aplicación React funcional!\n\nEn la próxima sesión, exploraremos los 'Hooks' de React más a fondo para manejar efectos secundarios y contextos."
    })

    # --- FIN DE LAS DIAPOSITIVAS ---

    os.makedirs(RUTA_PRESENTACIONES, exist_ok=True)

    nombre_archivo = 'presentacion_S12_s1_elementos_react.pptx'
    ruta_salida = os.path.join(RUTA_PRESENTACIONES, nombre_archivo)
    try:
        prs.save(ruta_salida)
        print(f"Presentación de la Semana 12 - Sesión 1 (JavaScript Avanzado) generada exitosamente en: {ruta_salida}")
    except Exception as e:
        print(f"Error al guardar la presentación: {e}")

def generar_presentacion_semana13_1_javascript():
    """
    Genera una presentación completa para la semana 13, sesión 1 del curso de 
    JavaScript Avanzado.
    Tema: Hooks en React: state, context, ref, y resource hooks.
    """
    print("Iniciando generación de presentación Semana 13 - Sesión 1 (JavaScript Avanzado - Hooks en React)")
    
    # Cargar la plantilla y los datos del curso
    ruta_plantilla = os.path.join(RUTA_PLANTILLAS, 'plantilla_2.pptx')
    ruta_json = os.path.join(RUTA_JSON, 'JavaScript.json')

    if not os.path.exists(ruta_plantilla):
        print(f"Error: No se encontró la plantilla en {ruta_plantilla}")
        return
    
    if not os.path.exists(ruta_json):
        print(f"Error: No se encontró el archivo JSON en {ruta_json}")
        return

    prs = Presentation(ruta_plantilla)
    datos_curso_completo = cargar_json(ruta_json) 
    if not datos_curso_completo or 'silabo' not in datos_curso_completo: 
        print("Error: No se pudieron cargar los datos del curso o estructura JSON inválida.")
        return
    datos_curso = datos_curso_completo['silabo']
    docente_nombre = "Ms. Johan Max A. López Heredia"

    # --- Diapositiva de Portada ---
    slide_portada = crear_diapositiva_con_marcadores(prs, "Layout_Portada_UTP_1")
    titulo_portada = "S13.s1 - Hooks en React: Superpoderes para tus Componentes" 
    reemplazar_marcadores(slide_portada, {
        "TITULO_SEMANA": titulo_portada,
        "NOMBRE_SEMANA": f"{datos_curso.get('curso', 'Curso Desconocido')} - {datos_curso.get('codigo', 'Código Desconocido')}",
        "DOCENTE_NOMBRE": docente_nombre
    })
    
    # --- Momento 1: Inicio ---
    slide_inicio = crear_diapositiva_con_marcadores(prs, "Layout_Inicio")
    reemplazar_marcadores(slide_inicio, {
        "TITULO_INICIO": "Más Allá de los Componentes Básicos",
        "ACTIVIDAD_ROMPEHIELO": (
            "Actividad inicial (5 minutos):\n\n"
            "Piensa en un superhéroe:\n\n"
            "- ¿Qué lo hace especial? (Sus poderes: volar, súper fuerza, etc.).\n"
            "- ¿Podría hacer lo mismo sin esos poderes?\n"
            "- ¿Cómo le ayudan esos poderes a resolver problemas complejos?\n\n"
            "Los Hooks en React son como 'superpoderes' que le damos a nuestros componentes funcionales para que puedan hacer cosas asombrosas que antes eran muy complicadas."
        )
    })
    
    # --- Momento 2: Utilidad ---
    slide_logro = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_Logro")
    logro_sesion = ("Al finalizar la sesión, el estudiante comprende el propósito de los Hooks en React y "
                    "diferencia el uso de los hooks básicos (useState, useEffect, useContext) para manejar el estado, "
                    "los efectos secundarios y el contexto en una aplicación web.")
    
    reemplazar_marcadores(slide_logro, {
        "ENCABEZADO_LOGRO": "Logro de Aprendizaje de la Sesión",
        "LOGRO_TEXTO": logro_sesion
    })
    
    slide_dudas = crear_diapositiva_con_marcadores(prs, "Layout_Utilidad_dudas")
    reemplazar_marcadores(slide_dudas, {
        "ENCABEZADO_DUDAS": "Dudas Frecuentes / Repaso",
        "DUDAS_TEXTO": (
            "- ¿Qué es un 'Hook'? ¿Por qué se llama así?\n"
            "- ¿Solo se usan en componentes de función?\n"
            "- ¿`useState` es el único Hook que existe?\n"
            "- ¿Cómo comparto información entre componentes que están muy separados?\n"
            "- ¿Qué es un 'efecto secundario' en React?"
        ),
        "ENCABEZADO_IMPORTANCIA_TEMA": "Importancia de los Hooks",
        "IMPORTANCIA_TEXTO": (
            "Aprender a usar Hooks es fundamental en el React moderno porque:\n\n"
            "• Permiten usar estado y otras características de React sin escribir clases.\n"
            "• Hacen el código más simple, reutilizable y fácil de entender.\n"
            "• Facilitan la separación de la lógica compleja de los componentes.\n"
            "• Son la forma estándar de escribir componentes en React hoy en día.\n"
            "• Permiten crear tus propios 'superpoderes' (custom hooks) para reutilizar lógica."
        )
    })
    
    slide_saberes = crear_diapositiva_con_marcadores(prs, "Layout_Saberes")
    reemplazar_marcadores(slide_saberes, {
        "TITULO_CONOCIMIENTOS": "Recordando los Componentes Funcionales",
        "DINAMICA_PREVIA": (
            "Recordemos brevemente la sesión anterior (5 minutos):\n\n"
            "1. ¿Qué es un componente funcional en React? (Una función de JavaScript que devuelve JSX).\n\n"
            "2. ¿Cómo le dábamos 'memoria' o estado a un componente funcional? (Con el Hook `useState`).\n\n"
            "3. ¿Qué son las `props`? (Datos que un componente padre pasa a un hijo).\n\n"
            "Hoy veremos que `useState` es solo el comienzo. Hay muchos otros 'poderes' que podemos añadir a nuestros componentes."
        )
    })
    
    # --- Momento 3: Transformación ---
    slide_desarrollo1 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo1, {
        "SUBTEMA_TITULO": "¿Qué son los Hooks?",
        "CONTENIDO_TEXTO": (
            "Los Hooks son funciones especiales que te permiten 'engancharte' (hook into) a las características de React desde componentes funcionales.\n\n"
            "Antes de los Hooks, si necesitabas estado o ciclo de vida, tenías que usar componentes de clase, que eran más verbosos.\n\n"
            "Reglas de los Hooks:\n"
            "1. Solo se pueden llamar en el nivel superior de un componente funcional (no dentro de bucles, condiciones o funciones anidadas).\n"
            "2. Solo se pueden llamar desde componentes funcionales de React (no desde funciones JavaScript normales).\n\n"
            "Todos los hooks empiezan con la palabra `use` (ej. `useState`, `useEffect`)."
        )
    })
    
    slide_desarrollo2 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo2, {
        "SUBTEMA_TITULO": "State Hook: `useState` (Repaso y Profundización)",
        "CONTENIDO_TEXTO": (
            "Ya lo conocemos, pero es el Hook más fundamental.\n\n"
            "• Propósito: Añadir estado local a un componente.\n"
            "• ¿Qué hace?: Declara una 'variable de estado'. React preservará esta variable entre renderizados.\n"
            "• Devuelve: Un par de valores: el estado actual y una función para actualizarlo.\n\n"
            "Ejemplo: `const [contador, setContador] = useState(0);`\n\n"
            "Es la base para hacer que nuestros componentes sean dinámicos e interactivos."
        )
    })

    slide_desarrollo3 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo3, {
        "SUBTEMA_TITULO": "Effect Hook: `useEffect`",
        "CONTENIDO_TEXTO": (
            "• Propósito: Permite ejecutar 'efectos secundarios' en componentes funcionales.\n\n"
            "¿Qué es un efecto secundario?\n"
            "Cualquier operación que afecte a algo fuera del componente, como:\n"
            "- Peticiones a una API (fetch de datos).\n"
            "- Suscripciones (ej. a un chat en tiempo real).\n"
            "- Manipulación directa del DOM.\n\n"
            "`useEffect` se ejecuta después de cada renderizado (por defecto), o cuando cambian ciertas dependencias que le especifiquemos."
        )
    })

    slide_desarrollo4 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo_Codigo")
    reemplazar_marcadores(slide_desarrollo4, {
        "SUBTEMA_TITULO_CODIGO": "Ejemplo Conceptual de `useEffect`",
        "EJEMPLO_CODIGO": (
            "import { useState, useEffect } from 'react';\n\n"
            "function TituloDinamico() {\n"
            "  const [nombre, setNombre] = useState('Usuario');\n\n"
            "  // Este efecto se ejecuta cada vez que 'nombre' cambia\n"
            "  useEffect(() => {\n"
            "    // Efecto secundario: cambiar el título del documento HTML\n"
            "    document.title = `Bienvenido, ${nombre}`;\n"
            "  }, [nombre]); // Array de dependencias\n\n"
            "  return (\n"
            "    <input \n"
            "      value={nombre}\n"
            "      onChange={(e) => setNombre(e.target.value)}\n"
            "    />\n"
            "  );\n"
            "}"
        )
    })

    slide_desarrollo5 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo5, {
        "SUBTEMA_TITULO": "Context Hook: `useContext`",
        "CONTENIDO_TEXTO": (
            "• Propósito: Compartir datos entre componentes sin tener que pasarlos manualmente a través de `props` por todos los niveles del árbol de componentes.\n\n"
            "Resuelve el problema del 'prop drilling' (perforación de props).\n\n"
            "Funciona con la API de Contexto de React:\n"
            "1. Se crea un `Context` (ej. `ThemeContext`).\n"
            "2. Se envuelve un árbol de componentes con un `Context.Provider` que provee un valor (ej. 'dark' o 'light').\n"
            "3. Cualquier componente hijo dentro de ese árbol puede 'consumir' ese valor usando el hook `useContext`."
        )
    })

    slide_desarrollo6 = crear_diapositiva_con_marcadores(prs, "Layout_Desarrollo")
    reemplazar_marcadores(slide_desarrollo6, {
        "SUBTEMA_TITULO": "Ref Hook: `useRef`",
        "CONTENIDO_TEXTO": (
            "• Propósito: Permite crear una referencia mutable que persiste durante todo el ciclo de vida del componente.\n\n"
            "Dos usos principales:\n"
            "1. Acceder a un elemento del DOM directamente: Para manejar el foco, medir su tamaño, etc.\n"
            "   `const miInputRef = useRef(null);`\n"
            "   `<input ref={miInputRef} />`\n\n"
            "2. Guardar cualquier valor mutable que no cause un nuevo renderizado cuando cambia (como una variable de instancia en una clase).\n"
            "   `const idIntervalo = useRef(null);`"
        )
    })
    
    # --- Momento 4: Práctica ---
    slide_practica = crear_diapositiva_con_marcadores(prs, "Layout_Practica")
    reemplazar_marcadores(slide_practica, {
        "TITULO_PRACTICA": "Práctica Conceptual: Aplicando Hooks",
        "INSTRUCCIONES_PRACTICA": (
            "Actividad en Equipos (Conceptual):\n\n"
            "Para cada uno de los siguientes escenarios, discutan qué Hook de React sería el más apropiado para resolver el problema y por qué.\n\n"
            "1. Escenario: Necesitas mostrar un contador de 'likes' en un botón. Cada vez que el usuario hace clic, el número debe aumentar en la pantalla.\n\n"
            "2. Escenario: Quieres que el título de la pestaña del navegador se actualice para mostrar el nombre del producto que el usuario está viendo.\n\n"
            "3. Escenario: Tienes un botón 'Modo Oscuro' en la cabecera de tu app, y quieres que todos los componentes (botones, textos, fondos) cambien de color cuando se activa, sin pasar la prop 'tema' a cada uno.\n\n"
            "4. Escenario: Quieres que al cargar una página de login, el cursor aparezca automáticamente en el campo de texto del email."
        ),
        "RECURSOS_NECESARIOS": (
            "• Papel y lápiz o editor de texto para tomar notas.\n"
            "• Tiempo estimado: 20 minutos."
        )
    })
    
    # --- Momento 5: Cierre ---
    slide_cierre = crear_diapositiva_con_marcadores(prs, "Layout_Cierre")
    reemplazar_marcadores(slide_cierre, {
        "TITULO_CIERRE": "Resumen: Los Superpoderes de React (Hooks)",
        "RESUMEN_IDEAS": (
            "En esta sesión hemos aprendido sobre los Hooks de React:\n\n"
            "• Son funciones que nos permiten 'engancharnos' a las características de React.\n"
            "• `useState`: Para añadir estado y hacer componentes dinámicos.\n"
            "• `useEffect`: Para manejar efectos secundarios como peticiones a APIs o manipulación del DOM.\n"
            "• `useContext`: Para compartir datos globalmente sin 'prop drilling'.\n"
            "• `useRef`: Para acceder a elementos del DOM o guardar valores mutables sin re-renderizar."
        ),
        "PREGUNTAS_METACOGNITIVAS": (
            "Reflexionemos:\n\n"
            "1. ¿Cómo crees que los Hooks simplifican el código en comparación con tener que escribir componentes de clase?\n"
            "2. ¿En qué situación de una aplicación real sería indispensable usar `useEffect`?\n"
            "3. ¿Qué problema resuelve `useContext` que sería muy tedioso de manejar solo con `props`?"
        )
    })
    
    slide_preguntas = crear_diapositiva_con_marcadores(prs, "Layout_Preguntas")
    reemplazar_marcadores(slide_preguntas, {
        "TITULO_PREGUNTAS": "¿Dudas sobre los Hooks de React?",
        "DUDAS_LISTADO": (
            "Ahora es el momento de resolver dudas sobre:\n\n"
            "• El propósito de cada uno de los Hooks vistos.\n"
            "• Las 'Reglas de los Hooks'.\n"
            "• El array de dependencias en `useEffect`.\n"
            "• La diferencia entre `useState` y `useRef`.\n"
            "• Casos de uso prácticos para `useContext`.\n\n"
            "¡Asegurémonos de entender estos conceptos clave!"
        )
    })
    
    slide_final = crear_diapositiva_con_marcadores(prs, "1_Layout_Preguntas")
    reemplazar_marcadores(slide_final, {
        "TITULO_AGRADECIMIENTO": "¡Felicidades por adquirir nuevos superpoderes en React!\n\nEn la próxima sesión, veremos cómo usar estos Hooks para conectarnos a servicios externos y APIs."
    })

    # --- FIN DE LAS DIAPOSITIVAS ---

    os.makedirs(RUTA_PRESENTACIONES, exist_ok=True)
    nombre_archivo = 'presentacion_S13_s1_hooks_react.pptx'
    ruta_salida = os.path.join(RUTA_PRESENTACIONES, nombre_archivo)
    try:
        prs.save(ruta_salida)
        print(f"Presentación de la Semana 13 - Sesión 1 (JavaScript Avanzado) generada exitosamente en: {ruta_salida}")
    except Exception as e:
        print(f"Error al guardar la presentación: {e}")
